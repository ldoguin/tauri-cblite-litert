"""
Generate workshop.pptx using the Couchbase Brand Template 2025.

Usage:
    python3 make_pptx.py

Requirements:
    pip install python-pptx

The template provides slide layouts, fonts (Open Sans), and the Couchbase
visual identity. This script adds slides using those layouts and injects
content via placeholders where possible, falling back to text boxes for
code blocks and diagrams.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pygments import lex
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.token import Token
from pygments.util import ClassNotFound
from pptx.oxml.ns import qn
from lxml import etree
import copy

TEMPLATE = "Couchbase Brand Template 2025.pptx"
OUTPUT   = "workshop.pptx"

# ── Couchbase brand colours (from theme2.xml "Couchbase Colors") ──────────────
# accent2=#EC1218 (red), accent3=#FC9C0C (orange), accent6=#BA0E12 (dark red)
CB_RED    = RGBColor(0xEC, 0x12, 0x18)   # accent2 — primary brand red
CB_DKRED  = RGBColor(0xBA, 0x0E, 0x12)   # accent6 — dark red (hover/active states)
CB_ORANGE = RGBColor(0xFC, 0x9C, 0x0C)   # accent3 — brand orange
CB_CREAM  = RGBColor(0xFF, 0xF0, 0xDB)   # accent4 — light cream (layout bg shapes)
CB_DARK   = RGBColor(0x00, 0x00, 0x00)   # dk1     — primary text (black)
CB_GREY   = RGBColor(0x59, 0x59, 0x59)   # dk2-ish — secondary text
CB_LGREY  = RGBColor(0xE5, 0xE5, 0xE5)   # lt2     — light grey background
CB_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # lt1
# Non-theme colours used for code and accents
# CB_TEAL / CB_GREEN removed — replaced by Couchbase brand colours below.
# Aliases kept so slide code needs no changes.
CB_TEAL   = CB_ORANGE                    # orange replaces teal for inline labels
CB_GREEN  = CB_DKRED                     # dark red replaces green for positive accents
CODE_BG   = RGBColor(0xF2, 0xF2, 0xF2)   # code block background
CODE_FG   = RGBColor(0x21, 0x21, 0x21)   # code text

FONT      = "Open Sans"

prs = Presentation(TEMPLATE)
SW = prs.slide_width    # 10"
SH = prs.slide_height   # 5.625"

# Remove all existing slides from the template (keep layouts/masters only).
# drop_rel() removes the slide part from the OPC package; removing the sldId
# element drops it from the presentation XML. Together they leave a clean slate.
_NS_R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
_sldIdLst = prs.slides._sldIdLst
for _sldId in list(_sldIdLst.sldId_lst):
    _rId = _sldId.get(_NS_R + 'id')
    prs.part.drop_rel(_rId)
    _sldIdLst.remove(_sldId)

# Use Master 1 (the rich Couchbase-branded master) for all slides.
# Master 0 has plain layouts with no visual identity.
# Master 1 layout indices (from template inspection):
#   0  'Title Slide 1'               — title slide with logo + decorative shapes
#   5  'SECTION_HEADER'              — red diagonal bg, ph idx=0 (title), idx=1 (subtitle)
#  17  'TITLE_ONLY'                  — content slide, ph idx=0 (title), 4 brand shapes
#  28  'Blank Layout'                — full canvas, ph idx=0 (hidden title), 3 brand shapes
_M1 = prs.slide_masters[1]

def new_slide(layout_idx):
    return prs.slides.add_slide(_M1.slide_layouts[layout_idx])

def _set_ph(slide, idx, text, size=None, bold=False, color=None,
            align=PP_ALIGN.LEFT, italic=False):
    """Write text into a placeholder by index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            p.clear()
            run = p.add_run()
            run.text = text
            if size:   run.font.size  = size
            if bold:   run.font.bold  = bold
            if italic: run.font.italic = italic
            if color:  run.font.color.rgb = color
            run.font.name = FONT
            return ph
    return None

def box(slide, l, t, w, h, fc=None, lc=None, lw=Pt(0)):
    """Add a filled rectangle."""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.width = lw
    if fc:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc
    else:
        sh.line.fill.background()
    return sh

def txb(slide, text, l, t, w, h, size=Pt(14), bold=False,
        color=CB_DARK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a free text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.bold   = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name   = FONT
    return tb

def _token_color(ttype):
    """Map a Pygments token type to an RGBColor for a light background."""
    t = ttype
    while t:
        if t in (Token.Keyword, Token.Keyword.Declaration,
                 Token.Keyword.Namespace, Token.Keyword.Reserved):
            return RGBColor(0xBA, 0x0E, 0x12)   # dark red — keywords
        if t is Token.Keyword.Type:
            return RGBColor(0x7B, 0x00, 0x8B)   # purple — types
        if t in (Token.Name.Function, Token.Name.Function.Magic):
            return RGBColor(0x00, 0x5C, 0xC8)   # blue — functions
        if t in (Token.Name.Builtin, Token.Name.Builtin.Pseudo):
            return RGBColor(0x00, 0x70, 0x70)   # teal — builtins
        if t in (Token.Name.Class, Token.Name.Decorator):
            return RGBColor(0x6F, 0x42, 0xC1)   # purple — classes
        if t is Token.Literal.String or str(t).startswith("Token.Literal.String"):
            return RGBColor(0x18, 0x7A, 0x1F)   # green — strings
        if t is Token.Literal.Number or str(t).startswith("Token.Literal.Number"):
            return RGBColor(0xC0, 0x50, 0x00)   # orange — numbers
        if t is Token.Comment or str(t).startswith("Token.Comment"):
            return RGBColor(0x6A, 0x6A, 0x6A)   # grey — comments
        if t in (Token.Operator, Token.Punctuation):
            return RGBColor(0x44, 0x44, 0x44)   # dark grey — operators
        t = t.parent
    return CODE_FG

_LANG_HINTS = [
    ("async fn ", "rust"), ("#[tauri", "rust"), ("use std::", "rust"),
    ("impl ", "rust"), ("let mut ", "rust"), ("-> Result<", "rust"),
    ("interface ", "typescript"), ("export ", "typescript"),
    ("import {", "typescript"), ("const ", "typescript"),
    ("await invoke", "typescript"), ("useCallback", "typescript"),
    ("def ", "python"), ("from pptx", "python"),
    ("SELECT ", "sql"), ("FROM _default", "sql"),
    ("$ cargo", "bash"), ("$ npm", "bash"), ("$ git", "bash"),
]

def _detect_lang(text: str) -> str:
    for marker, lang in _LANG_HINTS:
        if marker in text:
            return lang
    return "text"

def code_block(slide, text, l, t, w, h, size=Pt(9), lang=None):
    """Monospaced code block with Pygments syntax highlighting."""
    box(slide, l, t, w, h, fc=CODE_BG, lc=RGBColor(0xCC, 0xCC, 0xCC), lw=Pt(0.5))
    tb = slide.shapes.add_textbox(
        l + Inches(0.12), t + Inches(0.08),
        w - Inches(0.24), h - Inches(0.16))
    tf = tb.text_frame
    tf.word_wrap = False

    detected = lang or _detect_lang(text)
    try:
        lexer = get_lexer_by_name(detected, stripall=False)
        tokens = list(lex(text.rstrip("\n"), lexer))
    except ClassNotFound:
        tokens = [(Token.Text, text)]

    # Split token stream into lines
    lines: list[list[tuple]] = [[]]
    for ttype, value in tokens:
        parts = value.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                lines.append([])
            if part:
                lines[-1].append((ttype, part))

    first = True
    for line_tokens in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        if not line_tokens:
            r = p.add_run(); r.text = " "
            r.font.name = "Courier New"; r.font.size = size
        for ttype, value in line_tokens:
            r = p.add_run()
            r.text = value
            r.font.name  = "Courier New"
            r.font.size  = size
            r.font.color.rgb = _token_color(ttype)
            r.font.bold  = ttype in (Token.Keyword, Token.Keyword.Declaration,
                                      Token.Keyword.Namespace)

def label(slide, text, l, t, w, h, color=CB_TEAL, size=Pt(11)):
    """Small coloured label / badge."""
    box(slide, l, t, w, h, fc=color)
    txb(slide, text, l + Inches(0.08), t, w - Inches(0.08), h,
        size=size, bold=True, color=CB_WHITE, align=PP_ALIGN.LEFT)

def buls(slide, items, l, t, w, h, size=Pt(13), color=CB_DARK, bullet="\u25b8  "):
    """Bullet list."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = bullet + item
        r.font.size  = size
        r.font.color.rgb = color
        r.font.name  = FONT

def table(slide, headers, rows, l, t, w, h, hc=CB_RED, size=Pt(11)):
    """Simple table rendered as coloured boxes."""
    cw = w / len(headers)
    rh = h / (len(rows) + 1)
    for i, hdr in enumerate(headers):
        box(slide, l + cw*i, t, cw, rh, fc=hc)
        txb(slide, hdr, l + cw*i + Inches(0.06), t, cw, rh,
            size=size, bold=True, color=CB_WHITE)
    for r, row in enumerate(rows):
        ry = t + rh*(r+1)
        fc = CB_WHITE if r % 2 == 0 else RGBColor(0xF3,0xF3,0xF3)
        for c, cell in enumerate(row):
            box(slide, l + cw*c, ry, cw, rh, fc=fc,
                lc=RGBColor(0xCC,0xCC,0xCC), lw=Pt(0.25))
            txb(slide, cell, l + cw*c + Inches(0.06), ry, cw, rh,
                size=size-Pt(1), color=CB_DARK)

# ── Diagram helpers ───────────────────────────────────────────────────────────

def diag_box(slide, label, sublabel, l, t, w, h,
             fc=None, tc=CB_DARK, size=Pt(11)):
    """Filled rectangle with centred label (and optional smaller sublabel)."""
    fc = fc or RGBColor(0xF2, 0xF2, 0xF2)
    bg = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    bg.fill.solid(); bg.fill.fore_color.rgb = fc
    bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC); bg.line.width = Pt(0.75)
    tf = bg.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top  = tf.margin_bottom = Inches(0.04)
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = size; r.font.color.rgb = tc; r.font.bold = True; r.font.name = FONT
    if sublabel:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sublabel
        r2.font.size = Pt(size.pt - 2); r2.font.color.rgb = CB_GREY; r2.font.name = FONT
    return bg

def diag_arrow(slide, x1, y1, x2, y2, color=CB_RED, w=Pt(1.5)):
    """Straight connector arrow from (x1,y1) to (x2,y2)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    cx = abs(x2 - x1); cy = abs(y2 - y1)
    lx = min(x1, x2); ly = min(y1, y2)
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT
    conn.line.color.rgb = color; conn.line.width = w
    return conn

def diag_label(slide, text, l, t, w, h, color=CB_GREY, size=Pt(9), bold=False):
    """Small centred label (no background)."""
    txb(slide, text, l, t, w, h, size=size, color=color,
        align=PP_ALIGN.CENTER, bold=bold)

def part_tag(slide, text):
    """Small part label in top-right, below the title bar, in brand red."""
    txb(slide, text,
        SW - Inches(3.5), Inches(0.82), Inches(3.4), Inches(0.28),
        size=Pt(9), color=CB_RED, align=PP_ALIGN.RIGHT, italic=True)

# ── Slide factory functions ───────────────────────────────────────────────────
# Master 1 layout indices used:
#   0  Title Slide 1    — opening title slide
#   5  SECTION_HEADER   — part dividers (has idx=0 title + idx=1 subtitle)
#  17  TITLE_ONLY       — content slides (idx=0 title at 0.49", 0.14")
#  28  Blank Layout     — full-canvas slides

def title_slide(title, subtitle):
    """Master1 Layout 0 — opening title slide with Couchbase branding."""
    s = new_slide(0)
    _set_ph(s, 0, title,    size=Pt(32), bold=True,  color=CB_DARK, align=PP_ALIGN.LEFT)
    _set_ph(s, 1, subtitle, size=Pt(14), bold=False, color=CB_GREY, align=PP_ALIGN.LEFT)
    return s

def section_slide(title, subtitle=""):
    """Master1 Layout 5 — SECTION_HEADER with red diagonal background."""
    s = new_slide(5)
    _set_ph(s, 0, title,    size=Pt(28), bold=True, color=CB_DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        _set_ph(s, 1, subtitle, size=Pt(14), color=CB_GREY, align=PP_ALIGN.LEFT)
    return s

def content_slide(title, part_label=None):
    """Master1 Layout 17 — TITLE_ONLY with brand shapes; free canvas below title."""
    s = new_slide(17)
    _set_ph(s, 0, title, size=Pt(20), bold=True, color=CB_RED)
    return s

# Canvas constants — content area below the title bar (title ph ends ~0.79")
CY = Inches(1.0)   # content top y
CH = Inches(4.4)   # content height
CW = Inches(9.16)  # content width  (matches ph width in Master1 Layout 17)
CX = Inches(0.49)  # content left x (matches ph left in Master1 Layout 17)

# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = title_slide(
    "Build an Offline AI Desktop App",
    "Tauri  \u00b7  LiteRT  \u00b7  Couchbase Lite\n\nDeveloper Workshop  \u2014  61 slides  \u00b7  9 parts"
)
# Coloured tech badges
for i, (lbl, col) in enumerate([(  "Tauri v2", CB_TEAL),
                                  ("LiteRT-LM", CB_GREEN),
                                  ("Couchbase Lite", CB_RED)]):
    bx = Inches(0.34 + i * 3.1)
    box(s, bx, Inches(4.5), Inches(2.8), Inches(0.38), fc=col)
    txb(s, lbl, bx + Inches(0.08), Inches(4.5), Inches(2.7), Inches(0.38),
        size=Pt(13), bold=True, color=CB_WHITE, align=PP_ALIGN.CENTER)

# =============================================================================
# PART 1 — Introduction
# =============================================================================
section_slide("Part 1\nIntroduction & Problem Statement", "Slides 2\u20135")

# ── Slide 2 — The Problem ─────────────────────────────────────────────────────
s = content_slide("The Problem with Cloud AI", "Part 1 \u00b7 Introduction")
for bx, lbl, col in [(CX, "User query", CB_TEAL),
                      (Inches(3.8), "Cloud API", CB_RED),
                      (Inches(7.2), "Response",  CB_GREEN)]:
    box(s, bx, Inches(1.5), Inches(2.0), Inches(0.55), fc=col)
    txb(s, lbl, bx+Inches(0.05), Inches(1.5), Inches(1.9), Inches(0.55),
        size=Pt(13), bold=True, color=CB_WHITE, align=PP_ALIGN.CENTER)
for ax in [Inches(2.4), Inches(5.85)]:
    txb(s, "\u2500\u2500\u2500\u25ba", ax, Inches(1.45), Inches(1.3), Inches(0.55),
        size=Pt(16), color=CB_GREY, align=PP_ALIGN.CENTER)
txb(s, "\u26a0  Your data leaves the machine",
    Inches(3.0), Inches(2.2), Inches(4.0), Inches(0.4),
    size=Pt(12), color=CB_RED, bold=True, align=PP_ALIGN.CENTER)
buls(s, ["Privacy \u2014 medical, legal, financial data cannot leave the device",
         "Latency \u2014 200\u20132000 ms round-trip to a remote datacenter",
         "Cost \u2014 token pricing at scale (GPT-4o: ~$5 / M input tokens)",
         "Availability \u2014 no network = no AI"],
     CX, Inches(2.75), CW, Inches(2.5), size=Pt(14))

# ── Slide 3 — The Solution ────────────────────────────────────────────────────
s = content_slide("The Solution: On-Device AI", "Part 1 \u00b7 Introduction")
for bx, lbl, col in [(CX, "User query", CB_TEAL),
                      (Inches(3.8), "Local Model", CB_GREEN),
                      (Inches(7.2), "Response",    CB_TEAL)]:
    box(s, bx, Inches(1.5), Inches(2.0), Inches(0.55), fc=col)
    txb(s, lbl, bx+Inches(0.05), Inches(1.5), Inches(1.9), Inches(0.55),
        size=Pt(13), bold=True, color=CB_WHITE, align=PP_ALIGN.CENTER)
for ax in [Inches(2.4), Inches(5.85)]:
    txb(s, "\u2500\u2500\u2500\u25ba", ax, Inches(1.45), Inches(1.3), Inches(0.55),
        size=Pt(16), color=CB_GREY, align=PP_ALIGN.CENTER)
txb(s, "\u2713  Data never leaves the device",
    Inches(3.0), Inches(2.2), Inches(4.0), Inches(0.4),
    size=Pt(12), color=CB_GREEN, bold=True, align=PP_ALIGN.CENTER)
txb(s, "What we need:", CX, Inches(2.75), CW, Inches(0.35),
    size=Pt(14), bold=True, color=CB_TEAL)
buls(s, ["A runtime that executes quantized LLMs on CPU/GPU  \u2192  LiteRT",
         "A database that stores vectors and documents offline  \u2192  Couchbase Lite",
         "A shell that packages a web UI as a native desktop app  \u2192  Tauri"],
     CX, Inches(3.1), CW, Inches(2.0), size=Pt(14))

# ── Slide 4 — Architecture ────────────────────────────────────────────────────
s = content_slide("What We\u2019re Building")
# Stacked architecture diagram
_bw = Inches(5.5); _bx = CX; _gap = Inches(0.08)
_layers = [
    ("React UI  (Vite + TypeScript)",    "Chat · Knowledge · Model manager · Agents", RGBColor(0xFF,0xF0,0xDB)),
    ("Tauri v2 IPC bridge",              "JSON over postMessage",                      RGBColor(0xEC,0x12,0x18)),
    ("tauri-plugin-litert",              "LLM inference",                              RGBColor(0xF2,0xF2,0xF2)),
    ("tauri-plugin-cblite",              "Document + vector storage",                  RGBColor(0xF2,0xF2,0xF2)),
    ("Rust process  (reqwest · tokio)",  "Model download · HTTP · File I/O",           RGBColor(0xE5,0xE5,0xE5)),
]
_bh = Inches(0.72); _ty = CY
for i, (lbl, sub, fc) in enumerate(_layers):
    tc = CB_WHITE if fc == RGBColor(0xEC,0x12,0x18) else CB_DARK
    # plugins side by side
    if lbl.startswith("tauri-plugin-litert"):
        diag_box(s, "tauri-plugin-litert", "LLM inference",          _bx,                  _ty, Inches(2.65), _bh, fc=RGBColor(0xF2,0xF2,0xF2))
        diag_box(s, "tauri-plugin-cblite", "Document + vector store", _bx+Inches(2.65)+_gap, _ty, Inches(2.77), _bh, fc=RGBColor(0xF2,0xF2,0xF2))
    elif lbl.startswith("tauri-plugin-cblite"):
        pass  # drawn above
    else:
        diag_box(s, lbl, sub, _bx, _ty, _bw, _bh, fc=fc, tc=tc)
    if i < len(_layers) - 1 and not lbl.startswith("tauri-plugin-cblite"):
        _ty += _bh + _gap
buls(s, ["One-shot chat with local Gemma model",
         "Persistent history in Couchbase Lite",
         "RAG over PDFs, URLs, text",
         "Tool use (calculator, Wikipedia\u2026)",
         "Multi-agent routing"],
     Inches(6.1), CY + Inches(0.2), Inches(3.55), Inches(3.5), size=Pt(13))

# ── Slide 5 — Prerequisites ───────────────────────────────────────────────────
s = content_slide("Repository & Prerequisites", "Part 1 \u00b7 Introduction")
code_block(s, "git clone https://github.com/ldoguin/tauri-cblite-litert\ncd tauri-cblite-litert\npnpm install",
           CX, CY, CW, Inches(0.7), size=Pt(10))
table(s, ["Tool", "Version", "Purpose"],
      [["Rust",      ">= 1.77", "Tauri backend"],
       ["Node.js",   ">= 20",   "Frontend build"],
       ["pnpm",      ">= 9",    "Package manager"],
       ["Tauri CLI", "v2",      "cargo tauri dev"],
       ["patchelf",  "any",     "Linux .so RUNPATH fixup"]],
      CX, Inches(2.1), Inches(6.5), Inches(2.8), size=Pt(11))
txb(s, "Browser mode:", Inches(7.1), Inches(2.1), Inches(2.5), Inches(0.35),
    size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, "pnpm dev\n# http://localhost:1420",
           Inches(7.1), Inches(2.45), Inches(2.5), Inches(0.6), size=Pt(9))
txb(s, "Desktop app:", Inches(7.1), Inches(3.15), Inches(2.5), Inches(0.35),
    size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, "cargo tauri dev",
           Inches(7.1), Inches(3.5), Inches(2.5), Inches(0.4), size=Pt(9))

# =============================================================================
# PART 2 — Tauri
# =============================================================================
section_slide("Part 2\nTauri", "Slides 6\u201310")

# ── Slide 6 ───────────────────────────────────────────────────────────────────
s = content_slide("What is Tauri?")
_bw = Inches(4.3); _bx = CX; _bh = Inches(0.82); _gap = Inches(0.07); _ty = CY
for lbl, sub, fc in [
    ("WebView  (OS-native renderer)", "Your React / Vue / Svelte app",  RGBColor(0xFF,0xF0,0xDB)),
    ("Tauri IPC",                     "JSON over postMessage",           RGBColor(0xEC,0x12,0x18)),
    ("Rust process",                  "File system · HTTP · Native APIs",RGBColor(0xE5,0xE5,0xE5)),
]:
    tc = CB_WHITE if fc == RGBColor(0xEC,0x12,0x18) else CB_DARK
    diag_box(s, lbl, sub, _bx, _ty, _bw, _bh, fc=fc, tc=tc)
    _ty += _bh + _gap
table(s, ["", "Tauri", "Electron"],
      [["Binary size", "~10 MB",    "~150 MB"],
       ["Memory",      "~50 MB",    "~200 MB"],
       ["Renderer",    "OS WebView","Bundled Chromium"],
       ["Backend",     "Rust",      "Node.js"],
       ["Security",    "ACL / cmd", "Broad Node access"]],
      Inches(5.1), CY, Inches(4.56), Inches(3.2), size=Pt(11))

# ── Slide 7 ───────────────────────────────────────────────────────────────────
s = content_slide("Tauri Project Structure", "Part 2 \u00b7 Tauri")
code_block(s, (
    "tauri-cblite-litert/\n"
    "\u251c\u2500\u2500 src/                    # React frontend\n"
    "\u2502   \u251c\u2500\u2500 App.tsx\n"
    "\u2502   \u251c\u2500\u2500 hooks/useChat.ts    # Central state\n"
    "\u2502   \u2514\u2500\u2500 lib/                # Business logic\n"
    "\u251c\u2500\u2500 src-tauri/\n"
    "\u2502   \u251c\u2500\u2500 src/lib.rs          # Rust commands\n"
    "\u2502   \u251c\u2500\u2500 Cargo.toml\n"
    "\u2502   \u251c\u2500\u2500 tauri.conf.json\n"
    "\u2502   \u251c\u2500\u2500 capabilities/default.json  # ACL\n"
    "\u2502   \u2514\u2500\u2500 build.rs\n"
    "\u2514\u2500\u2500 packages/\n"
    "    \u251c\u2500\u2500 tauri-plugin-cblite/\n"
    "    \u2514\u2500\u2500 tauri-plugin-litert-api/"
), CX, CY, CW, Inches(3.8), size=Pt(10))

# ── Slide 8 ───────────────────────────────────────────────────────────────────
s = content_slide("Tauri IPC: Rust \u2194 TypeScript", "Part 2 \u00b7 Tauri")
txb(s, "Rust command:", CX, CY, Inches(4.6), Inches(0.3), size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, (
    "#[tauri::command]\nasync fn get_model_path(\n"
    "    app: AppHandle, file_name: String,\n"
    ") -> Result<Option<String>, String> {\n"
    "    let path = app.path()\n"
    "        .app_local_data_dir()?\n"
    "        .join(\"models\").join(&file_name);\n"
    "    Ok(path.exists().then(|| path.to_string_lossy().into()))\n"
    "}"
), CX, Inches(1.6), Inches(4.6), Inches(2.2), size=Pt(9))
txb(s, "TypeScript call:", Inches(5.1), CY, Inches(4.56), Inches(0.3), size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, (
    "import { invoke } from \"@tauri-apps/api/core\";\n\n"
    "const path = await invoke<string | null>(\n"
    "  \"get_model_path\",\n"
    "  { fileName: \"gemma3-1b-it-int4.litertlm\" }\n"
    ");"
), Inches(5.1), Inches(1.6), Inches(4.56), Inches(1.8), size=Pt(9))
txb(s, "Register in lib.rs:", CX, Inches(3.9), CW, Inches(0.3), size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, (
    "tauri::Builder::default()\n"
    "    .invoke_handler(tauri::generate_handler![get_model_path])\n"
    "    .run(tauri::generate_context!())"
), CX, Inches(4.2), CW, Inches(0.85), size=Pt(9))

# ── Slide 9 ───────────────────────────────────────────────────────────────────
s = content_slide("Tauri ACL: Security by Default", "Part 2 \u00b7 Tauri")
code_block(s, (
    "{\n  \"identifier\": \"default\",\n"
    "  \"platforms\": [\"linux\", \"macOS\", \"windows\"],\n"
    "  \"windows\": [\"main\"],\n"
    "  \"permissions\": [\n"
    "    \"core:default\",\n    \"os:default\",\n"
    "    \"cblite:default\",\n    \"litert:default\"\n"
    "  ]\n}"
), CX, CY, Inches(4.6), Inches(2.8), size=Pt(9))
txb(s, "cblite:default expands to:", Inches(5.1), CY, Inches(4.56), Inches(0.3),
    size=Pt(11), bold=True, color=CB_RED)
code_block(s, (
    "allow-open-database\nallow-close-database\n"
    "allow-get-document\nallow-save-document\n"
    "allow-execute-query\nallow-save-blob\n"
    "allow-get-blob-data\nallow-register-listener\n"
    "allow-remove-listener"
), Inches(5.1), Inches(1.6), Inches(4.56), Inches(2.5), size=Pt(9))
txb(s, "\u26a0  Without a capability entry, invoke() returns \"not found\"",
    CX, Inches(4.3), CW, Inches(0.4), size=Pt(12), color=CB_RED, bold=True)

# ── Slide 10 ──────────────────────────────────────────────────────────────────
s = content_slide("Isomorphic Pattern: Tauri + Web Fallback", "Part 2 \u00b7 Tauri")
code_block(s, (
    "export function isTauri(): boolean {\n"
    "  return typeof window !== \"undefined\"\n"
    "    && \"__TAURI_INTERNALS__\" in window;\n"
    "}\n\n"
    "export async function saveConversation(conv) {\n"
    "  const { id, ...body } = conv;\n"
    "  if (!isTauri()) {\n"
    "    webStore.set(\"conversations\", id, body); return;\n"
    "  }\n"
    "  const { saveDocument } = await import(\"tauri-plugin-cblite\");\n"
    "  await saveDocument(\"_default.conversations\", id, body);\n"
    "}"
), CX, CY, Inches(6.0), Inches(3.8), size=Pt(9))
buls(s, ["Fast iteration in the browser (no Rust compile)",
         "Same test suite covers both paths",
         "Graceful degradation if Tauri APIs unavailable",
         "Dynamic import() prevents browser evaluation of Tauri modules"],
     Inches(6.3), CY + Inches(0.3), Inches(3.36), Inches(3.0), size=Pt(12))

# =============================================================================
# PART 3 — AI Fundamentals
# =============================================================================
section_slide("Part 3\nAI Fundamentals", "Slides 11\u201323")

# ── Slide 11 ──────────────────────────────────────────────────────────────────
s = content_slide("Large Language Models in One Slide", "Part 3 \u00b7 AI Fundamentals")
code_block(s, '"The capital of France is"  ->  ["Paris": 0.94, "Lyon": 0.02, ...]',
           CX, CY, CW, Inches(0.5), size=Pt(10))
table(s, ["Concept", "What it means in practice"],
      [["Context window", "Max tokens the model can see at once (e.g. 8192)"],
       ["Quantization",   "FP32 -> INT4 = 8x smaller, ~5% quality loss"],
       ["Temperature",    "Randomness: 0 = deterministic, 1 = creative"],
       ["System prompt",  "Instructions prepended before the conversation"],
       ["Streaming",      "Emit tokens one by one, not waiting for full response"]],
      CX, Inches(1.9), CW, Inches(2.8), size=Pt(11))
txb(s, "Model: Gemma 3 1B INT4  (~700 MB)  \u00b7  ~30 tok/s CPU  \u00b7  ~130 tok/s GPU",
    CX, Inches(4.85), CW, Inches(0.4), size=Pt(11), color=CB_GREY, italic=True)

# ── Slide 12 ──────────────────────────────────────────────────────────────────
s = content_slide("Embeddings: Text \u2192 Vectors", "Part 3 \u00b7 AI Fundamentals")
code_block(s, (
    '"The dog barked"   ->  [0.12, -0.34, 0.89, ...]   (128 dims)\n'
    '"A canine howled"  ->  [0.11, -0.31, 0.91, ...]   (similar!)\n'
    '"The stock rose"   ->  [-0.45, 0.67, -0.12, ...]  (different)'
), CX, CY, CW, Inches(0.9), size=Pt(10))
code_block(s, "similarity = (A . B) / (|A| x |B|)   in [-1, 1]",
           CX, Inches(2.3), Inches(5.5), Inches(0.45), size=Pt(10))
txb(s, "Cosine similarity:", CX, Inches(2.1), Inches(3.5), Inches(0.3),
    size=Pt(11), bold=True, color=CB_TEAL)
for i, (lbl, pct, col) in enumerate([
    ('"dog" vs "canine"', 0.93, CB_GREEN),
    ('"dog" vs "puppy"',  0.81, CB_GREEN),
    ('"dog" vs "stock"',  0.12, CB_RED),
]):
    y = Inches(2.95) + Inches(0.5) * i
    txb(s, lbl, CX, y, Inches(2.5), Inches(0.4), size=Pt(11), color=CB_GREY)
    box(s, Inches(2.9), y + Inches(0.07), Inches(pct * 4.5), Inches(0.28), fc=col)
    txb(s, f"{pct:.0%}", Inches(2.9) + Inches(pct*4.5) + Inches(0.1), y,
        Inches(0.5), Inches(0.4), size=Pt(11), color=CB_DARK)
txb(s, "In this app: BERT-base-uncased  \u2192  128-dim embeddings  (~25 MB model)",
    CX, Inches(4.85), CW, Inches(0.4), size=Pt(11), color=CB_GREY, italic=True)

# ── Slide 13 ──────────────────────────────────────────────────────────────────
s = content_slide("The BERT Tokenizer (WordPiece)", "Part 3 \u00b7 AI Fundamentals")
code_block(s, '"unaffable"  ->  ["un", "##aff", "##able"]  ->  [2512, 4593, 3085]',
           CX, CY, CW, Inches(0.5), size=Pt(10))
txb(s, "Special tokens:", CX, Inches(1.95), Inches(3.5), Inches(0.3),
    size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, "[CLS] token1 token2 ... [SEP] [PAD] [PAD]\n 101                    102    0     0\n\nThe [CLS] output vector is used as the sentence embedding.",
           CX, Inches(2.25), Inches(5.5), Inches(1.2), size=Pt(9))
txb(s, "TypeScript implementation:", CX, Inches(3.55), Inches(5.5), Inches(0.3),
    size=Pt(11), bold=True, color=CB_TEAL)
code_block(s, "const tokens = wordpieceTokenize(text, vocab);\nconst inputIds = [CLS_ID, ...tokens.slice(0, 126), SEP_ID];\nwhile (inputIds.length < 128) inputIds.push(PAD_ID);",
           CX, Inches(3.85), Inches(5.5), Inches(1.1), size=Pt(9))
buls(s, ["Full WordPiece tokenizer in ~100 lines TypeScript",
         "30,522-token bert-base-uncased vocabulary",
         "Vocab cached in Cache API after first load"],
     Inches(5.8), Inches(2.0), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 14 ──────────────────────────────────────────────────────────────────
s = content_slide("RAG: Retrieval-Augmented Generation")
# Two-row pipeline diagram
_bh = Inches(0.55); _bw_sm = Inches(1.55); _gap = Inches(0.18); _arr = Inches(0.22)
# Row 1: Ingest
txb(s, "INGEST  (offline)", CX, CY, Inches(2.2), Inches(0.3), size=Pt(9), color=CB_GREY, bold=True)
_row1 = [("Document", RGBColor(0xFF,0xF0,0xDB)), ("Chunk", RGBColor(0xF2,0xF2,0xF2)),
         ("Embed", RGBColor(0xF2,0xF2,0xF2)), ("Store in DB", RGBColor(0xE5,0xE5,0xE5))]
_x = CX; _y1 = CY + Inches(0.32)
for j,(lbl,fc) in enumerate(_row1):
    diag_box(s, lbl, "", _x, _y1, _bw_sm, _bh, fc=fc, size=Pt(10))
    _x += _bw_sm
    if j < len(_row1)-1:
        diag_label(s, "\u2192", _x, _y1, _arr, _bh, color=CB_RED, size=Pt(14), bold=True)
        _x += _arr
# Row 2: Query
txb(s, "QUERY  (online)", CX, _y1+_bh+Inches(0.25), Inches(2.2), Inches(0.3), size=Pt(9), color=CB_GREY, bold=True)
_row2 = [("User query", RGBColor(0xFF,0xF0,0xDB)), ("Embed", RGBColor(0xF2,0xF2,0xF2)),
         ("Search DB", RGBColor(0xF2,0xF2,0xF2)), ("Top-K chunks", RGBColor(0xE5,0xE5,0xE5))]
_x = CX; _y2 = _y1 + _bh + Inches(0.57)
for j,(lbl,fc) in enumerate(_row2):
    diag_box(s, lbl, "", _x, _y2, _bw_sm, _bh, fc=fc, size=Pt(10))
    _x += _bw_sm
    if j < len(_row2)-1:
        diag_label(s, "\u2192", _x, _y2, _arr, _bh, color=CB_RED, size=Pt(14), bold=True)
        _x += _arr
# Down arrow + LLM box
_mid_x = CX + Inches(3*(_bw_sm+_arr)) + _bw_sm/2 - Inches(0.9)
diag_label(s, "\u2193", _mid_x, _y2+_bh, Inches(1.8), Inches(0.35), color=CB_RED, size=Pt(14), bold=True)
diag_box(s, "Inject into LLM prompt", "", _mid_x - Inches(0.1), _y2+_bh+Inches(0.35), Inches(2.0), _bh, fc=RGBColor(0xEC,0x12,0x18), tc=CB_WHITE, size=Pt(10))
diag_label(s, "\u2193", _mid_x, _y2+_bh*2+Inches(0.35), Inches(1.8), Inches(0.35), color=CB_RED, size=Pt(14), bold=True)
diag_box(s, "LLM generates answer", "", _mid_x - Inches(0.1), _y2+_bh*2+Inches(0.7), Inches(2.0), _bh, fc=RGBColor(0xE5,0xE5,0xE5), size=Pt(10))
buls(s, ["No fine-tuning \u2014 just an embedding model + vector store",
         "Grounds responses in your private documents",
         "Chunking: overlapping windows preserve cross-boundary context"],
     Inches(5.9), CY + Inches(0.2), Inches(3.76), Inches(3.5), size=Pt(12))

# ── Slide 15 ──────────────────────────────────────────────────────────────────
s = content_slide("Hybrid Search: Vector + BM25 + RRF", "Part 3 \u00b7 AI Fundamentals")
code_block(s, (
    "BM25 score(q,d) = SUM IDF(t) * tf(t,d)*(k1+1) / (tf(t,d) + k1*(1-b+b*|d|/avgdl))\n"
    "                  k1=1.5  b=0.75\n\n"
    "RRF(d) = SUM weight_i / (k + rank_i(d))    k=60"
), CX, CY, CW, Inches(1.2), size=Pt(9))
code_block(s, (
    "Vector results:  [doc3, doc1, doc7, doc2, ...]\n"
    "BM25 results:    [doc1, doc5, doc3, doc8, ...]\n"
    "                       |  RRF fusion  (bm25Weight=0.3)\n"
    "Merged results:  [doc1, doc3, doc5, doc7, ...]"
), CX, Inches(2.65), CW, Inches(1.3), size=Pt(9))
buls(s, ["Pure vector search misses exact keyword matches",
         "Pure BM25 misses semantic similarity",
         "RRF k=60 dampens rank differences at the top",
         "Tune bm25Weight: higher for legal/medical, lower for semantic queries"],
     CX, Inches(4.1), CW, Inches(1.5), size=Pt(13))

# ── Slide 16 ──────────────────────────────────────────────────────────────────
s = content_slide("The ReAct Loop: LLMs That Use Tools", "Part 3 \u00b7 AI Fundamentals")
code_block(s, (
    "User: What is the GDP of France divided by its population?\n\n"
    "LLM:  I need the GDP and population of France.\n"
    "      <tool_call>{\"name\":\"wikipedia\",\"args\":{\"query\":\"France GDP\"}}</tool_call>\n\n"
    "Tool: France GDP is $3.1 trillion (2023)\n"
    "      <tool_call>{\"name\":\"wikipedia\",\"args\":{\"query\":\"France population\"}}</tool_call>\n\n"
    "Tool: France population is 68 million\n\n"
    "LLM:  GDP per capita = $3.1T / 68M = $45,588"
), CX, CY, Inches(6.5), Inches(3.6), size=Pt(9))
buls(s, ["Max 5 iterations per message",
         "Duplicate tool call \u2192 loop guard \u2192 break",
         "<tool_call> tags stripped from displayed response"],
     Inches(6.8), CY + Inches(0.3), Inches(2.86), Inches(2.5), size=Pt(12))

# ── Slide 17 ──────────────────────────────────────────────────────────────────
s = content_slide("Agent Routing")
_bw = Inches(3.2); _bh = Inches(0.52); _bx = CX + Inches(0.3); _gap = Inches(0.12); _arr = Inches(0.28)
_steps = [
    ("User message",                    RGBColor(0xFF,0xF0,0xDB)),
    ("Router LLM call",                 RGBColor(0xEC,0x12,0x18)),
    ('{ "agent": "Support" }',          RGBColor(0xF2,0xF2,0xF2)),
    ("Support agent  +  tools",         RGBColor(0xF2,0xF2,0xF2)),
    ("Response",                        RGBColor(0xE5,0xE5,0xE5)),
]
_ty = CY
for i,(lbl,fc) in enumerate(_steps):
    tc = CB_WHITE if fc == RGBColor(0xEC,0x12,0x18) else CB_DARK
    diag_box(s, lbl, "", _bx, _ty, _bw, _bh, fc=fc, tc=tc, size=Pt(11))
    _ty += _bh
    if i < len(_steps)-1:
        diag_label(s, "\u2193", _bx, _ty, _bw, _arr, color=CB_RED, size=Pt(13), bold=True)
        _ty += _arr
# Side label for router input
txb(s, "\u2190  agent list + descriptions", _bx + _bw + Inches(0.1),
    CY + _bh + _arr/2, Inches(2.5), Inches(0.4), size=Pt(9), color=CB_GREY, italic=True)
buls(s, ["Every message goes through the router first",
         "Router uses the same LLM (~50 tokens in, ~10 out)",
         "Adds ~200 ms latency \u2014 acceptable for interactive chat",
         "Falls back to default system prompt if routing fails"],
     Inches(5.0), CY + Inches(0.3), Inches(4.66), Inches(3.0), size=Pt(13))

# ── Slide 18 ── Transformer Architecture ─────────────────────────────────────
s = content_slide("Transformer Architecture (How LLMs Work)", "Part 3 \u00b7 AI Fundamentals")
# Left: vertical block diagram of one transformer layer
_bw = Inches(2.4); _bh = Inches(0.46); _farr = Inches(0.16)
_fx = CX + Inches(0.1); _fy = CY
_layers = [
    ("Input Embeddings",        CB_LGREY,  CB_DARK),
    ("Positional Encoding",     CB_LGREY,  CB_DARK),
    ("Multi-Head Self-Attention",CB_RED,   CB_WHITE),
    ("Add & LayerNorm",         CB_CREAM,  CB_DARK),
    ("Feed-Forward Network",    CB_CREAM,  CB_DARK),
    ("Add & LayerNorm",         CB_CREAM,  CB_DARK),
    ("\u00d7 N layers",          CB_LGREY,  CB_GREY),
    ("Output Logits + Softmax", CB_DKRED,  CB_WHITE),
]
for i, (lbl, fc, tc) in enumerate(_layers):
    diag_box(s, lbl, "", _fx, _fy, _bw, _bh, fc=fc, tc=tc, size=Pt(9))
    _fy += _bh
    if i < len(_layers) - 1:
        diag_label(s, "\u2193", _fx, _fy, _bw, _farr, color=CB_RED, size=Pt(11), bold=True)
        _fy += _farr
# Right: key concepts
buls(s, [
    "Self-Attention: each token attends to all others \u2014 captures long-range dependencies",
    "Multi-Head: H parallel attention heads learn different relationship types",
    "Positional encoding: sine/cosine or learned \u2014 gives order to the sequence",
    "FFN: two linear layers with GELU activation \u2014 stores factual knowledge",
    "Residual connections (Add): prevent vanishing gradients in deep stacks",
    "Gemma 3 1B: 18 layers, 1152 hidden dim, 4 attention heads, 8192 ctx",
], Inches(3.1), CY + Inches(0.1), Inches(6.56), Inches(4.5), size=Pt(12))

# ── Slide 19 ── Attention Mechanism ──────────────────────────────────────────
s = content_slide("Self-Attention: The Core Mechanism", "Part 3 \u00b7 AI Fundamentals")
code_block(s, (
    "# For each token, compute Query, Key, Value projections\n"
    "Q = X @ W_Q    # (seq_len, d_k)\n"
    "K = X @ W_K    # (seq_len, d_k)\n"
    "V = X @ W_V    # (seq_len, d_v)\n\n"
    "# Scaled dot-product attention\n"
    "scores = (Q @ K.T) / sqrt(d_k)   # (seq_len, seq_len)\n"
    "weights = softmax(scores)          # each row sums to 1\n"
    "output  = weights @ V              # weighted sum of values\n\n"
    "# Causal mask (decoder): upper triangle = -inf so future tokens are invisible\n"
    "scores = scores + causal_mask"
), CX, CY, Inches(6.0), Inches(2.8), size=Pt(9))
buls(s, [
    "d_k = head_dim (e.g. 64) \u2014 sqrt(d_k) prevents vanishing gradients in softmax",
    "Each attention weight = how much token i should attend to token j",
    "Encoder (BERT): bidirectional \u2014 every token sees every other token",
    "Decoder (Gemma): causal mask \u2014 token i only sees tokens 0\u2026i",
    "KV-Cache: store K and V from previous tokens, only compute new token's Q",
    "Flash Attention: fused CUDA kernel \u2014 avoids materialising full (seq, seq) matrix",
], CX, Inches(3.95), CW, Inches(1.9), size=Pt(12))

# ── Slide 20 ── Embedding Models vs LLMs ─────────────────────────────────────
s = content_slide("Embedding Models vs Generative LLMs", "Part 3 \u00b7 AI Fundamentals")
table(s, ["", "Embedding model (BERT)", "Generative LLM (Gemma)"],
      [["Architecture",  "Encoder-only transformer",      "Decoder-only transformer"],
       ["Output",        "Dense vector (128\u2013768 dims)",    "Next-token probability distribution"],
       ["Training",      "MLM + NSP self-supervised",     "Causal language modelling"],
       ["Context",       "Bidirectional (sees full input)","Left-to-right (causal mask)"],
       ["Size",          "~25 MB (BERT-base INT8)",        "~700 MB (Gemma 3 1B INT4)"],
       ["Latency",       "~5 ms / 128 tokens (CPU)",      "~30\u2013130 tok/s (CPU/GPU)"],
       ["Use in app",    "Chunk embeddings for RAG",       "Chat, routing, tool use"]],
      CX, CY, CW, Inches(3.5), size=Pt(10))
buls(s, [
    "Never use a generative LLM to produce embeddings \u2014 its output space is not metric",
    "Sentence-BERT (SBERT) adds a pooling head to BERT for sentence-level similarity",
    "This app uses bert-base-uncased: 30k vocab, 12 layers, 768 hidden \u2192 128 projected",
], CX, Inches(4.65), CW, Inches(1.2), size=Pt(12))

# ── Slide 21 ── Why BERT Instead of a Dedicated Embedding Model ──────────────
s = content_slide("Why BERT Instead of a Dedicated Embedding Model?", "Part 3 \u00b7 AI Fundamentals")
txb(s, "Short answer: on-device size, LiteRT support, and good-enough quality for domain-specific RAG.",
    CX, CY, CW, Inches(0.4), size=Pt(12), color=CB_DARK, italic=True)
table(s, ["Criterion", "BERT-base-uncased (this app)", "all-MiniLM-L6-v2", "text-embedding-3-small"],
      [["Model size",      "~25 MB (INT8 .tflite)",        "~23 MB",              "Cloud API \u2014 no local file"],
       ["Output dims",     "768 \u2192 128 (projected)",        "384",                 "1536"],
       ["LiteRT support",  "\u2713 official .tflite",           "\u2713 convertible",        "\u2717 cloud only"],
       ["Runs fully offline", "\u2713",                         "\u2713",                   "\u2717 requires internet"],
       ["Latency (CPU)",   "~5 ms / 128 tokens",           "~4 ms / 128 tokens",  "~80\u2013200 ms (network)"],
       ["Quality (MTEB)",  "~38 avg",                      "~56 avg",             "~62 avg"],
       ["Token limit",     "128 tokens (hard)",            "256 tokens",          "8191 tokens"],
       ["Licence",         "Apache 2.0",                   "Apache 2.0",          "Proprietary / paid"]],
      CX, Inches(1.5), CW, Inches(3.5), size=Pt(9))
buls(s, [
    "BERT is not the best embedding model \u2014 all-MiniLM scores ~18 points higher on MTEB benchmarks",
    "It is used here because it has an official, well-tested .tflite export and runs on LiteRT without conversion work",
    "The 128-token limit is a real constraint: chunks must stay short, which is why DEFAULT_CHUNK_SIZE is 300 chars",
    "For production: swap in all-MiniLM-L6-v2 (convert to .tflite) for better cross-domain retrieval quality",
    "Cloud embedding APIs (OpenAI, Cohere) give the best quality but break the offline-first guarantee",
], CX, Inches(5.15), CW, Inches(1.2), size=Pt(11))

# ── Slide 23 ── Pooling & Projection ─────────────────────────────────────────
s = content_slide("From Token Vectors to Sentence Embeddings", "Part 3 \u00b7 AI Fundamentals")
txb(s, "Problem: BERT outputs one 768-dim vector per token. We need one vector per chunk.",
    CX, CY, CW, Inches(0.4), size=Pt(12), color=CB_DARK)
table(s, ["Pooling strategy", "How", "When to use"],
      [["[CLS] token",      "Take the first output vector (trained for classification)", "Classification, retrieval (BERT default)"],
       ["Mean pooling",     "Average all non-[PAD] token vectors",                       "Sentence similarity (SBERT)"],
       ["Max pooling",      "Element-wise max across token vectors",                     "Captures salient features"],
       ["Weighted mean",    "Weight by attention scores",                                "Rare \u2014 complex, marginal gain"]],
      CX, Inches(1.55), CW, Inches(2.2), size=Pt(10))
txb(s, "Projection layer (this app):", CX, Inches(3.85), Inches(4.0), Inches(0.35),
    size=Pt(12), bold=True, color=CB_ORANGE)
code_block(s, (
    "// After [CLS] pooling: 768-dim vector\n"
    "// Linear projection W \u2208 R^{768\u00d7128} reduces to 128 dims\n"
    "embedding_128 = cls_vector @ W_proj   // (128,)\n"
    "embedding_128 = l2_normalize(embedding_128)\n\n"
    "// L2 normalisation: ||v|| = 1  \u21d2  cosine_sim(a,b) = dot(a,b)\n"
    "// Stored as REAL[] in CouchbaseLite \u2014 128 floats = 512 bytes per chunk"
), CX, Inches(4.2), CW, Inches(1.9), size=Pt(9))

# ── Slide 22 ── Vector Dimensions & Distance Metrics ─────────────────────────
s = content_slide("Vector Dimensions & Distance Metrics", "Part 3 \u00b7 AI Fundamentals")
table(s, ["Metric", "Formula", "Range", "Best for"],
      [["Cosine similarity", "A\u00b7B / (|A||B|)",         "[\u22121, 1]",  "Normalised embeddings (this app)"],
       ["Dot product",       "A\u00b7B",                    "(\u2212\u221e,\u221e)", "When magnitude encodes relevance"],
       ["Euclidean (L2)",    "\u221a\u03a3(a\u1d62\u2212b\u1d62)\u00b2",  "[0, \u221e)",  "Image embeddings, k-NN"],
       ["Manhattan (L1)",    "\u03a3|a\u1d62\u2212b\u1d62|",             "[0, \u221e)",  "Sparse vectors, anomaly detection"]],
      CX, CY, CW, Inches(2.3), size=Pt(11))
txb(s, "Dimensionality trade-offs:", CX, Inches(3.4), Inches(4.0), Inches(0.35),
    size=Pt(12), bold=True, color=CB_ORANGE)
table(s, ["Dims", "Model example", "Storage/chunk", "Quality"],
      [["128",  "BERT projected (this app)",  "512 B",  "Good for short text retrieval"],
       ["384",  "all-MiniLM-L6-v2",           "1.5 KB", "Better cross-domain"],
       ["768",  "BERT-base, RoBERTa",          "3 KB",   "Strong general purpose"],
       ["1536", "text-embedding-ada-002",      "6 KB",   "State-of-art (cloud only)"]],
      CX, Inches(3.75), CW, Inches(2.0), size=Pt(10))
txb(s, "Higher dims \u2260 always better \u2014 128 dims is sufficient for domain-specific RAG with short chunks",
    CX, Inches(5.85), CW, Inches(0.35), size=Pt(11), color=CB_GREY, italic=True)

# =============================================================================
# PART 4 — Couchbase Lite
# =============================================================================
section_slide("Part 4\nCouchbase Lite", "Slides 24\u201328")

# ── Slide 24 ──────────────────────────────────────────────────────────────────
s = content_slide("Couchbase Lite: Embedded NoSQL", "Part 4 \u00b7 Couchbase Lite")
table(s, ["Feature", "Detail"],
      [["Storage",    "SQLite-backed, MVCC, ACID transactions"],
       ["Data model", "JSON documents in named collections"],
       ["Query",      "SQL++ (superset of SQL for JSON)"],
       ["Sync",       "Sync Gateway / Capella App Services"],
       ["Blobs",      "Binary attachments stored outside JSON"],
       ["Listeners",  "Real-time change notifications"],
       ["Platforms",  "iOS, Android, macOS, Windows, Linux, .NET, JVM"]],
      CX, CY, CW, Inches(3.5), size=Pt(12))
txb(s, "In this app: tauri-plugin-cblite wraps the C SDK via Rust FFI",
    CX, Inches(4.85), CW, Inches(0.4), size=Pt(11), color=CB_GREY, italic=True)

# ── Slide 25 ──────────────────────────────────────────────────────────────────
s = content_slide("Collections & Document Model")
# Database schema as a visual hierarchy
_dbw = Inches(1.6); _dbh = Inches(0.48); _cw = Inches(3.5); _ch = Inches(0.42)
_dbx = CX; _dby = CY
diag_box(s, "app.cblite2", "_default scope", _dbx, _dby, _dbw, _dbh,
         fc=RGBColor(0xEC,0x12,0x18), tc=CB_WHITE, size=Pt(10))
_cols = [
    ("conversations",  "id, title, agentId, createdAt"),
    ("messages",       "id, conversationId, role, content"),
    ("chunks",         "id, docId, text, embedding[128]"),
    ("documents",      "id, title, mimeType, blobKey"),
    ("app-config",     "schemaVersion"),
]
_cy2 = _dby; _cx2 = _dbx + _dbw + Inches(0.35)
for lbl, sub in _cols:
    diag_box(s, lbl, sub, _cx2, _cy2, _cw, _ch,
             fc=RGBColor(0xF2,0xF2,0xF2), size=Pt(9))
    # connector line
    diag_label(s, "\u2500", _dbx+_dbw, _cy2+_ch/2-Inches(0.1), Inches(0.35), Inches(0.3),
               color=CB_GREY, size=Pt(10))
    _cy2 += _ch + Inches(0.06)
code_block(s, (
    "// Save\nawait saveDocument(\"_default.messages\", msgId, {\n"
    "  conversationId, role: \"user\", content: text,\n"
    "  createdAt: Date.now()\n"
    "});\n\n"
    "// Query\nconst rows = await executeQuery(\n"
    "  \"SELECT id, role, content FROM _default.messages\"\n"
    "  + \" WHERE conversationId = $cid ORDER BY createdAt\",\n"
    "  { cid: conversationId }\n"
    ");"
), CX, Inches(3.1), CW, Inches(2.1), size=Pt(9))

# ── Slide 26 ──────────────────────────────────────────────────────────────────
s = content_slide("SQL++ Queries", "Part 4 \u00b7 Couchbase Lite")
code_block(s, (
    "-- All conversations, newest first\nSELECT id, title, agentId, createdAt\n"
    "FROM _default.conversations\nORDER BY createdAt DESC;\n\n"
    "-- Messages for a conversation\nSELECT id, role, content, toolCalls, toolResults\n"
    "FROM _default.messages\nWHERE conversationId = $cid\nORDER BY createdAt;\n\n"
    "-- Chunks with embeddings for a document\nSELECT id, text, embedding\n"
    "FROM _default.chunks\nWHERE docId = $docId;"
), CX, CY, CW, Inches(3.8), size=Pt(9))
buls(s, ["Parameters use $name syntax (prevents injection)",
         "ARRAY_CONTAINS(), ARRAY_LENGTH() for JSON arrays",
         "No JOINs across collections in CBL \u2014 denormalise or query separately"],
     CX, Inches(4.5), CW, Inches(1.1), size=Pt(13))

# ── Slide 27 ──────────────────────────────────────────────────────────────────
s = content_slide("Blobs: Storing Binary Data", "Part 4 \u00b7 Couchbase Lite")
code_block(s, (
    "// Store a PDF\nconst blobKey = await saveBlob(\n"
    "  \"application/pdf\",\n"
    "  Array.from(new Uint8Array(pdfBuffer))\n"
    ");\nawait saveDocument(\"_default.documents\", docId, {\n"
    "  title, mimeType: \"application/pdf\", blobKey\n"
    "});\n\n"
    "// Retrieve\nconst bytes = await getBlobData(blobKey);\nconst blob = new Blob([new Uint8Array(bytes)]);"
), CX, CY, Inches(5.5), Inches(3.5), size=Pt(9))
buls(s, ["Blob stored outside the document JSON",
         "Deduplicated by content hash",
         "blobKey is a stable reference across saves",
         "Sync Gateway replicates blobs separately (lazy)"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 28 ──────────────────────────────────────────────────────────────────
s = content_slide("Schema Migrations", "Part 4 \u00b7 Couchbase Lite")
code_block(s, (
    "async function runMigrations() {\n"
    "  const doc = await getDocument(\"_default\", \"app-config\")\n"
    "    .catch(e => e?.message?.includes(\"not found\") ? null : Promise.reject(e));\n"
    "  const version = doc?.schemaVersion ?? 0;\n\n"
    "  if (version < 1) {\n"
    "    // Create default agent\n"
    "    await saveDocument(\"_default.agents\", \"default\", { name: \"Assistant\" });\n"
    "  }\n"
    "  if (version < 2) {\n"
    "    // Add createdAt to existing conversations\n"
    "    const rows = await executeQuery(\"SELECT id FROM _default.conversations\");\n"
    "    for (const { id } of rows)\n"
    "      await saveDocument(\"_default.conversations\", id, { createdAt: 0 });\n"
    "  }\n"
    "  await saveDocument(\"_default\", \"app-config\", { schemaVersion: CURRENT_VERSION });\n"
    "}"
), CX, CY, CW, Inches(4.1), size=Pt(8.5))
txb(s, "Key: catch \"not found\" \u2014 CBL throws instead of returning null for missing docs",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_RED, bold=True)

# =============================================================================
# PART 5 — LiteRT
# =============================================================================
section_slide("Part 5\nLiteRT", "Slides 29\u201345")

# ── Slide 29 ──────────────────────────────────────────────────────────────────
s = content_slide("What is LiteRT?", "Part 5 \u00b7 LiteRT")
table(s, ["", "LiteRT (TFLite)", "llama.cpp", "ONNX Runtime"],
      [["Format",    ".litertlm / .tflite", "GGUF",        ".onnx"],
       ["Delegates", "GPU, NNAPI, CoreML",  "CUDA, Metal", "CUDA, DirectML"],
       ["Quantize",  "INT4, INT8, FP16",    "Q4_K_M etc.", "INT8, FP16"],
       ["Mobile",    "\u2713 first-class",  "\u2713",       "\u2713"],
       ["Desktop",   "\u2713",              "\u2713",       "\u2713"]],
      CX, CY, CW, Inches(3.0), size=Pt(11))
txb(s, "LiteRT-LM = LiteRT + GenAI API (KV-cache, sampling, streaming)",
    CX, Inches(4.3), CW, Inches(0.4), size=Pt(12), color=CB_ORANGE, bold=True)
txb(s, "Gemma 3 1B INT4 \u2014 official .litertlm from ai.google.dev",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_GREY, italic=True)

# ── Slide 30 ──────────────────────────────────────────────────────────────────
s = content_slide("Model Download: Rust Streaming", "Part 5 \u00b7 LiteRT")
code_block(s, (
    "#[tauri::command]\nasync fn download_model(\n"
    "    app: AppHandle, url: String, file_name: String,\n"
    "    on_event: Channel<DownloadEvent>,\n"
    ") -> Result<(), String> {\n"
    "    let dest = app.path().app_local_data_dir()?\n"
    "        .join(\"models\").join(&file_name);\n"
    "    let mut resp = reqwest::get(&url).await?;\n"
    "    let total = resp.content_length().unwrap_or(0);\n"
    "    let mut file = tokio::fs::File::create(&dest).await?;\n"
    "    let mut downloaded = 0u64;\n"
    "    while let Some(chunk) = resp.chunk().await? {\n"
    "        file.write_all(&chunk).await?;\n"
    "        downloaded += chunk.len() as u64;\n"
    "        on_event.send(DownloadEvent::Progress {\n"
    "            downloaded, total\n"
    "        })?;\n"
    "    }\n"
    "    on_event.send(DownloadEvent::Finished)?;\n"
    "    Ok(())\n"
    "}"
), CX, CY, CW, Inches(4.2), size=Pt(8.5))
txb(s, "Channel<T> = Tauri v2 streaming IPC \u2014 no polling needed",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_ORANGE, bold=True)

# ── Slide 31 ──────────────────────────────────────────────────────────────────
s = content_slide("Model Download: TypeScript Side", "Part 5 \u00b7 LiteRT")
code_block(s, (
    "import { Channel, invoke } from \"@tauri-apps/api/core\";\n\n"
    "const channel = new Channel<DownloadEvent>();\n"
    "channel.onmessage = (event) => {\n"
    "  if (event.type === \"progress\") {\n"
    "    const pct = event.downloaded / event.total;\n"
    "    setProgress(pct);\n"
    "  } else if (event.type === \"finished\") {\n"
    "    setDownloading(false);\n"
    "  }\n"
    "};\n\n"
    "await invoke(\"download_model\", {\n"
    "  url: MODEL_URL,\n"
    "  fileName: \"gemma3-1b-it-int4.litertlm\",\n"
    "  onEvent: channel,\n"
    "});"
), CX, CY, Inches(5.5), Inches(3.8), size=Pt(9))
buls(s, ["Channel is a typed event stream",
         "onEvent maps to on_event in Rust (camelCase \u2194 snake_case)",
         "Cancel: invoke(\"cancel_download\") sets an AtomicBool",
         "Progress persists across app restarts via modelCache.ts"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.8), size=Pt(12))

# ── Slide 32 ──────────────────────────────────────────────────────────────────
s = content_slide("Running Inference", "Part 5 \u00b7 LiteRT")
code_block(s, (
    "// TypeScript\nconst session = await invoke<string>(\"create_session\", {\n"
    "  modelPath,\n"
    "  systemPrompt: agent.systemPrompt,\n"
    "});\n\n"
    "const stream = new Channel<InferenceEvent>();\nstream.onmessage = (ev) => {\n"
    "  if (ev.type === \"token\") appendToken(ev.text);\n"
    "  else if (ev.type === \"done\") finalise();\n"
    "};\n\n"
    "await invoke(\"run_inference\", {\n"
    "  sessionId: session,\n"
    "  messages: history,   // [{role, content}]\n"
    "  onEvent: stream,\n"
    "});"
), CX, CY, Inches(5.5), Inches(3.8), size=Pt(9))
buls(s, ["Session holds KV-cache \u2014 reuse across turns",
         "History formatted as ChatML before passing to model",
         "Abort: invoke(\"abort_inference\", { sessionId })"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.0), size=Pt(12))

# ── Slide 33 ──────────────────────────────────────────────────────────────────
s = content_slide("Embedding Inference", "Part 5 \u00b7 LiteRT")
code_block(s, (
    "// Rust: run BERT embedding model\n#[tauri::command]\nasync fn run_embedding(\n"
    "    state: State<'_, LiteRTState>,\n"
    "    text: String,\n"
    ") -> Result<Vec<f32>, String> {\n"
    "    let interpreter = state.embedding_interpreter.lock().await;\n"
    "    let input = tokenize_bert(&text)?;\n"
    "    interpreter.set_input(0, &input)?;\n"
    "    interpreter.invoke()?;\n"
    "    let output = interpreter.get_output(0)?;\n"
    "    Ok(l2_normalize(output))\n"
    "}"
), CX, CY, Inches(5.5), Inches(3.0), size=Pt(9))
buls(s, ["BERT-base-uncased: 12 layers, 768 hidden, 128 output after pooling",
         "L2-normalise so cosine similarity = dot product",
         "~5 ms per chunk on CPU \u2014 fast enough for real-time ingest",
         "Web fallback: ONNX Runtime Web in a Worker"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 34 ── LiteRT Model Formats ─────────────────────────────────────────
s = content_slide("LiteRT Model Formats", "Part 5 \u00b7 LiteRT")
table(s, ["Format", "Contents", "Use case"],
      [[".tflite",    "FlatBuffer graph + weights",          "Classic TFLite ops (vision, audio)"],
       [".litertlm",  "FlatBuffer + GenAI metadata sidecar", "LLMs with KV-cache & sampling config"],
       [".task",      "Bundled model + pre/post-processing",  "MediaPipe Tasks (object detect, pose)"],
       ["SavedModel", "TF graph + variables directory",       "Training checkpoint \u2014 convert before deploy"]],
      CX, CY, CW, Inches(2.6), size=Pt(11))
txb(s, "Conversion pipeline:", CX, Inches(3.8), Inches(3.0), Inches(0.3),
    size=Pt(11), bold=True, color=CB_ORANGE)
code_block(s, (
    "# Keras / PyTorch  \u2192  SavedModel  \u2192  .tflite\n"
    "converter = tf.lite.TFLiteConverter.from_saved_model(saved_model_dir)\n"
    "converter.optimizations = [tf.lite.Optimize.DEFAULT]   # INT8 PTQ\n"
    "converter.target_spec.supported_types = [tf.float16]   # or INT4 via AI Edge\n"
    "tflite_model = converter.convert()\n\n"
    "# LLM-specific: use AI Edge Torch\n"
    "# pip install ai-edge-torch\n"
    "ai_edge_torch.convert(model, sample_inputs).export('gemma.litertlm')"
), CX, Inches(4.1), CW, Inches(2.1), size=Pt(8.5))

# ── Slide 35 ── Quantization Deep Dive ───────────────────────────────────────
s = content_slide("Quantization: Making Models Fit on Device", "Part 5 \u00b7 LiteRT")
table(s, ["Scheme", "Bits/weight", "Size vs FP32", "Quality loss", "Speed gain"],
      [["FP32",       "32", "1\u00d7  baseline",  "none",    "1\u00d7"],
       ["FP16",       "16", "0.5\u00d7",         "< 0.1%",  "1.5\u00d7 (GPU)"],
       ["INT8 (PTQ)", "8",  "0.25\u00d7",        "~1\u20132%", "2\u20134\u00d7 (CPU)"],
       ["INT4 (QAT)", "4",  "0.125\u00d7",       "~3\u20135%", "4\u20136\u00d7 (CPU)"],
       ["INT4 mixed", "4/8","0.15\u00d7",        "~1\u20132%", "3\u20135\u00d7"]],
      CX, CY, CW, Inches(2.8), size=Pt(11))
buls(s, [
    "PTQ = Post-Training Quantization \u2014 no retraining, calibration dataset only",
    "QAT = Quantization-Aware Training \u2014 fine-tune with fake quantization nodes",
    "Mixed precision: keep attention layers FP16, FFN layers INT4",
    "Gemma 3 1B INT4: 700 MB vs 4 GB FP32 \u2014 fits in phone RAM",
    "Rule of thumb: INT8 for accuracy-critical tasks, INT4 for chat/generation",
], CX, Inches(3.95), CW, Inches(1.8), size=Pt(12))

# ── Slide 36 ── Hardware Delegates ───────────────────────────────────────────
s = content_slide("Hardware Delegates: Choosing the Right Backend", "Part 5 \u00b7 LiteRT")
table(s, ["Delegate", "Platform", "Best for", "Fallback"],
      [["GPU Delegate",  "Android / iOS / Desktop", "LLM inference, matrix ops",    "CPU"],
       ["NNAPI",         "Android 8.1+",            "DSP / NPU on Qualcomm, MediaTek", "GPU \u2192 CPU"],
       ["CoreML",        "iOS / macOS",             "Apple Neural Engine (ANE)",    "CPU"],
       ["XNNPack",       "All (default CPU)",       "Optimised SIMD on ARM / x86",  "\u2014"],
       ["Hexagon DSP",   "Qualcomm Snapdragon",     "Ultra-low power inference",    "NNAPI"]],
      CX, CY, CW, Inches(2.8), size=Pt(11))
code_block(s, (
    "// LiteRT-LM selects delegate automatically via InferenceOptions\n"
    "const options: InferenceOptions = {\n"
    "  preferredBackend: 'gpu',   // 'gpu' | 'cpu' | 'nnapi'\n"
    "  numThreads: 4,             // CPU thread count when GPU unavailable\n"
    "};\n"
    "// Rust side: litert_lm::Session::new(model_path, options)"
), CX, Inches(3.95), CW, Inches(1.5), size=Pt(9))
txb(s, "GPU delegate gives 3\u20134\u00d7 speedup on Adreno / Mali GPUs vs XNNPack CPU",
    CX, Inches(5.55), CW, Inches(0.35), size=Pt(11), color=CB_ORANGE, bold=True)

# ── Slide 37 ── GenAI API & KV-Cache ─────────────────────────────────────────
s = content_slide("GenAI API: KV-Cache & Sampling", "Part 5 \u00b7 LiteRT")
txb(s, "What is a KV-Cache?", CX, CY, Inches(4.0), Inches(0.35),
    size=Pt(13), bold=True, color=CB_ORANGE)
buls(s, [
    "Transformers compute Key and Value matrices for every token in context",
    "Without cache: re-compute all K/V on every new token \u2014 O(n\u00b2) cost",
    "With KV-cache: store K/V from previous tokens, only compute new token \u2014 O(n)",
    "LiteRT-LM Session holds the cache across turns \u2014 reuse the same session object",
], CX, Inches(1.45), CW, Inches(1.6), size=Pt(12))
txb(s, "Sampling parameters:", CX, Inches(3.15), Inches(4.0), Inches(0.35),
    size=Pt(13), bold=True, color=CB_ORANGE)
table(s, ["Parameter", "Effect", "Typical value"],
      [["temperature",  "Scales logits before softmax \u2014 higher = more random",  "0.7\u20131.0"],
       ["top_k",        "Sample only from top-K most likely tokens",              "40"],
       ["top_p",        "Nucleus sampling: smallest set with cumulative p prob",  "0.95"],
       ["max_tokens",   "Hard stop on output length",                             "512\u20132048"],
       ["stop_tokens",  "EOS token IDs that terminate generation",                "model-specific"]],
      CX, Inches(3.5), CW, Inches(2.0), size=Pt(10))

# ── Slide 38 ── LiteRT Performance Benchmarks ─────────────────────────────────
s = content_slide("LiteRT Performance: Real Numbers", "Part 5 \u00b7 LiteRT")
table(s, ["Device", "Model", "Backend", "Prefill (tok/s)", "Decode (tok/s)"],
      [["Pixel 8 Pro",        "Gemma 3 1B INT4", "GPU (Adreno 740)", "~220", "~130"],
       ["Pixel 8 Pro",        "Gemma 3 1B INT4", "CPU (XNNPack)",    "~55",  "~30"],
       ["Samsung S24",        "Gemma 3 1B INT4", "GPU (Xclipse 940)","~180", "~110"],
       ["MacBook M3",         "Gemma 3 1B INT4", "CPU (XNNPack)",    "~400", "~250"],
       ["Desktop RTX 4070",   "Gemma 3 1B INT4", "GPU delegate",     "~900", "~600"],
       ["BERT-base (embed)",  "128 tokens",       "CPU",              "\u2014",  "~5 ms/call"]],
      CX, CY, CW, Inches(3.2), size=Pt(10))
buls(s, [
    "Prefill = processing the prompt (parallel) \u2014 Decode = generating tokens (sequential)",
    "First-token latency dominated by prefill \u2014 keep prompts short for responsiveness",
    "GPU delegate requires model to be loaded once; subsequent calls reuse weights",
    "BERT embedding: 5 ms/call \u2014 ingest 1000 chunks in ~5 s on CPU",
], CX, Inches(4.35), CW, Inches(1.5), size=Pt(12))

# ── Slide 39 ── LiteRT in the Tauri Plugin ────────────────────────────────────
s = content_slide("tauri-plugin-litert: Plugin Architecture", "Part 5 \u00b7 LiteRT")
# Two-column: left = plugin structure, right = command list
txb(s, "Plugin exposes Tauri commands:", CX, CY, Inches(4.5), Inches(0.35),
    size=Pt(12), bold=True, color=CB_ORANGE)
code_block(s, (
    "// src-tauri/capabilities/default.json\n"
    "\"permissions\": [\n"
    "  \"litert:allow-load-lm-model\",\n"
    "  \"litert:allow-run-inference\",\n"
    "  \"litert:allow-abort-inference\",\n"
    "  \"litert:allow-run-embedding\",\n"
    "  \"litert:allow-download-model\",\n"
    "  \"litert:allow-cancel-download\"\n"
    "]"
), CX, Inches(1.45), Inches(4.4), Inches(2.3), size=Pt(9))
txb(s, "TypeScript API surface:", Inches(5.0), CY, Inches(4.66), Inches(0.35),
    size=Pt(12), bold=True, color=CB_ORANGE)
code_block(s, (
    "import {\n"
    "  loadLmModel,      // load .litertlm into session\n"
    "  runInference,     // stream tokens via Channel\n"
    "  abortInference,   // cancel mid-stream\n"
    "  runEmbedding,     // BERT \u2192 float32[]\n"
    "  downloadModel,    // streaming download + progress\n"
    "  cancelDownload,   // abort in-flight download\n"
    "} from 'tauri-plugin-litert-api';"
), Inches(5.0), Inches(1.45), Inches(4.66), Inches(2.3), size=Pt(9))
buls(s, [
    "Plugin state: Arc<Mutex<Session>> shared across commands",
    "Embedding model loaded separately from LLM session",
    "All file I/O uses Tauri path resolver \u2014 sandboxed per platform",
    "Web fallback: same API shape, ONNX Runtime Web under the hood",
], CX, Inches(3.9), CW, Inches(1.8), size=Pt(12))

# ── Slide 40 ── TFLite History & the LiteRT Rename ───────────────────────────
s = content_slide("From TensorFlow Lite to LiteRT: The History", "Part 5 \u00b7 LiteRT")
# Timeline boxes
_bw = Inches(1.55); _bh = Inches(0.52); _farr = Inches(0.18)
_fx = CX; _fy = CY
_timeline = [
    ("2017",  "TensorFlow Lite announced at Google I/O",                CB_LGREY, CB_DARK),
    ("2019",  "TFLite 1.0 ships \u2014 Android + iOS, GPU delegate",          CB_LGREY, CB_DARK),
    ("2021",  "MediaPipe Solutions built on TFLite tasks",               CB_CREAM, CB_DARK),
    ("2023",  "GenAI API added \u2014 LLM inference, KV-cache, streaming",    CB_CREAM, CB_DARK),
    ("2024",  "TFLite rebranded \u2192 LiteRT under Google AI Edge umbrella", CB_RED,   CB_WHITE),
    ("2024+", "LiteRT-LM: dedicated LLM runtime (.litertlm format)",    CB_DKRED, CB_WHITE),
]
for i, (yr, desc, fc, tc) in enumerate(_timeline):
    diag_box(s, yr, "", _fx, _fy, _bw, _bh, fc=fc, tc=tc, size=Pt(10))
    txb(s, desc, _fx + _bw + Inches(0.14), _fy, Inches(7.0), _bh,
        size=Pt(10), color=CB_DARK)
    _fy += _bh
    if i < len(_timeline) - 1:
        diag_label(s, "\u2193", _fx, _fy, _bw, _farr, color=CB_RED, size=Pt(11), bold=True)
        _fy += _farr
buls(s, [
    "The .tflite format is unchanged \u2014 all existing models still work",
    "LiteRT is a rename + governance change, not a rewrite",
    "Google AI Edge = LiteRT + MediaPipe + AI Edge Torch (conversion toolchain)",
], CX, Inches(5.5), CW, Inches(0.9), size=Pt(11))

# ── Slide 41 ── MediaPipe & LiteRT: What\u2019s the Difference ─────────────────────
s = content_slide("MediaPipe vs LiteRT: Roles & Relationship", "Part 5 \u00b7 LiteRT")
table(s, ["", "LiteRT (core runtime)", "MediaPipe Solutions"],
      [["What it is",   "Low-level inference engine",              "High-level ML pipeline framework"],
       ["Abstraction",  "Runs a .tflite / .litertlm graph",       "Bundles model + pre/post-processing"],
       ["API level",    "Tensor in \u2192 tensor out",                  "Image in \u2192 structured result out"],
       ["Model format", ".tflite, .litertlm",                     ".task (wraps .tflite internally)"],
       ["Use cases",    "Custom models, LLMs, embeddings",        "Object detection, pose, face, text"],
       ["Who uses it",  "ML engineers, app developers",           "App developers, no ML expertise needed"],
       ["Relationship", "MediaPipe Tasks run on top of LiteRT",   "MediaPipe is a consumer of LiteRT"]],
      CX, CY, CW, Inches(3.6), size=Pt(10))
buls(s, [
    "MediaPipe Tasks = LiteRT + pre-built pipelines + .task bundle format",
    "If you need object detection or pose estimation: use MediaPipe Tasks",
    "If you need a custom model or LLM: use LiteRT directly (as this app does)",
    "Both ship in the same Google AI Edge SDK \u2014 not mutually exclusive",
], CX, Inches(4.75), CW, Inches(1.5), size=Pt(12))

# ── Slide 42 ── On-Device Model Management Strategies ────────────────────────
s = content_slide("On-Device Model Management: Strategies", "Part 5 \u00b7 LiteRT")
table(s, ["Strategy", "How", "Pros", "Cons"],
      [["Bundle in app",    "Ship model inside the binary / assets",
        "No download, always available",       "Large install size, hard to update"],
       ["Download on first run", "Fetch from CDN after install",
        "Small initial install",               "Requires internet, user waits"],
       ["Download on demand",    "Fetch when feature first used",
        "Only download if needed",             "Latency at feature activation"],
       ["Background prefetch",   "Download silently after install",
        "Ready before user needs it",          "Uses data without explicit consent"],
       ["Model hub / OTA",       "Server pushes new model versions",
        "Update without app store release",    "Complex infra, version management"]],
      CX, CY, CW, Inches(3.0), size=Pt(10))
buls(s, [
    "This app: download on demand \u2014 user explicitly triggers, progress shown via Channel",
    "Model stored in app_local_data_dir() \u2014 survives app updates, excluded from backup",
    "Integrity: verify SHA-256 after download before loading into LiteRT",
    "Multiple models: store in models/ subdirectory, keyed by filename",
    "Hugging Face Hub, ai.google.dev, and Kaggle Models are common CDN sources",
], CX, Inches(4.15), CW, Inches(1.7), size=Pt(12))

# ── Slide 43 ── On-Device Model Storage & Lifecycle ──────────────────────────
s = content_slide("Model Storage & Lifecycle on Each Platform", "Part 5 \u00b7 LiteRT")
table(s, ["Platform", "Storage path", "Persists across updates?", "Backup?"],
      [["Android",  "getFilesDir() / app_local_data_dir()",  "Yes (internal storage)",  "No (excluded)"],
       ["iOS",      "Application Support/",                   "Yes",                     "iCloud optional"],
       ["macOS",    "~/Library/Application Support/<bundle>", "Yes",                     "Time Machine"],
       ["Windows",  "%APPDATA%\\<bundle>\\",                  "Yes",                     "No default"],
       ["Linux",    "~/.local/share/<bundle>/",               "Yes",                     "No default"]],
      CX, CY, CW, Inches(2.6), size=Pt(10))
code_block(s, (
    "// Tauri resolves the correct path per platform automatically\n"
    "const modelDir = await appLocalDataDir();   // e.g. ~/Library/Application Support/com.app\n"
    "const modelPath = await join(modelDir, 'models', 'gemma3-1b-it-int4.litertlm');\n\n"
    "// Check if already downloaded before showing download UI\n"
    "const exists = await invoke<boolean>('model_exists', { path: modelPath });\n"
    "if (!exists) showDownloadPrompt();"
), CX, Inches(3.75), CW, Inches(1.9), size=Pt(9))
buls(s, [
    "Never store models in the app bundle directory \u2014 it may be read-only on some platforms",
    "Use Tauri's path API, not hardcoded paths \u2014 avoids /sdcard vs /storage/emulated/0 bugs",
], CX, Inches(5.75), CW, Inches(0.6), size=Pt(11))

# ── Slide 44 ── On-Device Runtime Comparison (Deep Dive) ─────────────────────
s = content_slide("On-Device Inference Runtimes: Full Comparison", "Part 5 \u00b7 LiteRT")
table(s, ["", "LiteRT-LM", "llama.cpp", "ONNX Runtime", "MLC-LLM", "Core ML"],
      [["Format",       ".litertlm",    "GGUF",         ".onnx",        "MLC weights",  ".mlpackage"],
       ["LLM support",  "\u2713 GenAI API","\u2713 native",  "via GenAI ext","native",       "\u2713 (iOS 18+)"],
       ["Mobile",       "\u2713 first",  "\u2713",         "\u2713",           "\u2713 Android", "iOS/macOS only"],
       ["GPU backend",  "Delegate API", "CUDA/Metal",   "DirectML/CUDA","Vulkan/Metal",  "ANE/GPU"],
       ["Streaming",    "\u2713 Channel", "\u2713 callback","\u2713 callback",  "\u2713",          "\u2713"],
       ["Rust binding", "\u2713 official","\u2713 llama-cpp-rs","ort crate",  "limited",      "\u2717"],
       ["Tauri plugin", "\u2713 this app","\u2717 DIY",     "\u2717 DIY",       "\u2717 DIY",      "\u2717 macOS only"],
       ["License",      "Apache 2.0",   "MIT",          "MIT",          "Apache 2.0",   "proprietary"]],
      CX, CY, CW, Inches(4.0), size=Pt(9))
buls(s, [
    "llama.cpp: best ecosystem for GGUF models (Mistral, Llama, Phi) \u2014 no official Tauri plugin",
    "ONNX Runtime: widest model compatibility, strong for non-LLM tasks (BERT, Whisper)",
    "MLC-LLM: WebGPU + Vulkan focus, good for browser + Android, less mature Rust story",
    "Core ML: Apple-only, best ANE utilisation on M-series \u2014 not cross-platform",
    "LiteRT: best choice when targeting Android + iOS + Desktop from one codebase",
], CX, Inches(5.15), CW, Inches(1.2), size=Pt(11))

# ── Slide 45 ── How to Choose a Runtime ──────────────────────────────────────
s = content_slide("How to Choose an On-Device Runtime", "Part 5 \u00b7 LiteRT")
# Decision table: question → answer → runtime(s) that fit
table(s, ["Your situation", "Reach for"],
      [
       ["You need Android + iOS + Desktop from one codebase",
        "LiteRT \u2014 first-class mobile delegates, single format"],
       ["Your model is already in GGUF (Llama, Mistral, Phi, Qwen \u2026)",
        "llama.cpp \u2014 widest GGUF ecosystem, mature quantisation"],
       ["You need non-LLM tasks: BERT, Whisper, YOLO, Stable Diffusion",
        "ONNX Runtime \u2014 broadest op coverage, strong tooling"],
       ["You target Apple devices only and want best ANE performance",
        "Core ML \u2014 deepest Apple Silicon integration"],
       ["You need WebGPU / browser + native from one model",
        "MLC-LLM \u2014 Vulkan + WebGPU, same weights on web and device"],
       ["You want a ready-made Tauri plugin with no Rust FFI work",
        "LiteRT (this repo) \u2014 plugin already written"],
       ["You need the largest model selection and community",
        "llama.cpp \u2014 most models, most quantisation recipes on HF Hub"],
       ["You need to run the same model in a Python training loop and on device",
        "ONNX Runtime \u2014 export once, run everywhere"],
      ],
      CX, CY, CW, Inches(4.1), size=Pt(10))
buls(s, [
    "No runtime wins on every axis \u2014 the right choice depends on your model source, target platforms, and team",
    "Runtimes are not mutually exclusive: use ONNX Runtime for BERT embeddings and LiteRT for the LLM in the same app",
    "Benchmark on your target device with your actual model \u2014 published numbers vary widely across hardware",
], CX, Inches(5.25), CW, Inches(1.1), size=Pt(11))

# =============================================================================
# PART 6 — RAG Pipeline
# =============================================================================
section_slide("Part 6\nRAG Pipeline", "Slides 46\u201350")

# ── Slide 46 ──────────────────────────────────────────────────────────────────
s = content_slide("Document Ingest Pipeline", "Part 6 \u00b7 RAG Pipeline")
code_block(s, (
    "async function ingestDocument(file: File) {\n"
    "  // 1. Extract text\n"
    "  const text = await extractText(file);          // PDF / TXT / MD\n\n"
    "  // 2. Store raw document + blob\n"
    "  const blobKey = await saveBlob(file.type, bytes);\n"
    "  const docId = await saveDocument(\"_default.documents\", uuid(), {\n"
    "    title: file.name, mimeType: file.type, blobKey\n"
    "  });\n\n"
    "  // 3. Chunk\n"
    "  const chunks = chunkText(text, { size: 512, overlap: 64 });\n\n"
    "  // 4. Embed + store\n"
    "  for (const chunk of chunks) {\n"
    "    const embedding = await embed(chunk);\n"
    "    await saveDocument(\"_default.chunks\", uuid(), {\n"
    "      docId, text: chunk, embedding\n"
    "    });\n"
    "  }\n"
    "}"
), CX, CY, CW, Inches(4.2), size=Pt(8.5))

# ── Slide 47 ──────────────────────────────────────────────────────────────────
s = content_slide("Chunking Strategy", "Part 6 \u00b7 RAG Pipeline")
code_block(s, (
    "function chunkText(text: string, opts = { size: 512, overlap: 64 }) {\n"
    "  const words = text.split(/\\s+/);\n"
    "  const chunks: string[] = [];\n"
    "  let i = 0;\n"
    "  while (i < words.length) {\n"
    "    chunks.push(words.slice(i, i + opts.size).join(\" \"));\n"
    "    i += opts.size - opts.overlap;\n"
    "  }\n"
    "  return chunks;\n"
    "}"
), CX, CY, Inches(5.5), Inches(2.5), size=Pt(9))
table(s, ["Parameter", "Value", "Effect"],
      [["size",    "512 words", "Fits BERT 512-token limit"],
       ["overlap", "64 words",  "Preserves cross-boundary context"],
       ["strategy","word-split","Language-agnostic, fast"]],
      CX, Inches(3.8), Inches(5.5), Inches(1.5), size=Pt(11))
buls(s, ["Semantic chunking (split on paragraphs/headings) gives better results",
         "Smaller chunks = more precise retrieval, more DB rows",
         "Larger chunks = more context per result, fewer rows"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 48 ──────────────────────────────────────────────────────────────────
s = content_slide("Vector Search Implementation", "Part 6 \u00b7 RAG Pipeline")
code_block(s, (
    "async function vectorSearch(queryEmbedding: number[], topK = 5) {\n"
    "  const rows = await executeQuery(\n"
    "    \"SELECT id, text, embedding FROM _default.chunks\"\n"
    "  );\n"
    "  return rows\n"
    "    .map(r => ({\n"
    "      ...r,\n"
    "      score: cosineSimilarity(queryEmbedding, r.embedding)\n"
    "    }))\n"
    "    .sort((a, b) => b.score - a.score)\n"
    "    .slice(0, topK);\n"
    "}\n\n"
    "function cosineSimilarity(a: number[], b: number[]) {\n"
    "  let dot = 0, na = 0, nb = 0;\n"
    "  for (let i = 0; i < a.length; i++) {\n"
    "    dot += a[i] * b[i]; na += a[i]**2; nb += b[i]**2;\n"
    "  }\n"
    "  return dot / (Math.sqrt(na) * Math.sqrt(nb));\n"
    "}"
), CX, CY, CW, Inches(4.1), size=Pt(8.5))
txb(s, "Note: full table scan \u2014 acceptable for <10k chunks; add HNSW index for larger corpora",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_RED, italic=True)

# ── Slide 49 ──────────────────────────────────────────────────────────────────
s = content_slide("BM25 + RRF Fusion", "Part 6 \u00b7 RAG Pipeline")
code_block(s, (
    "async function hybridSearch(query: string, topK = 5) {\n"
    "  const qEmbed = await embed(query);\n"
    "  const [vecResults, bm25Results] = await Promise.all([\n"
    "    vectorSearch(qEmbed, topK * 2),\n"
    "    bm25Search(query, topK * 2),\n"
    "  ]);\n"
    "  return rrfFusion(vecResults, bm25Results, {\n"
    "    k: 60, bm25Weight: 0.3, topK\n"
    "  });\n"
    "}\n\n"
    "function rrfFusion(vecR, bm25R, { k, bm25Weight, topK }) {\n"
    "  const scores = new Map<string, number>();\n"
    "  vecR.forEach((r, i) => scores.set(r.id,\n"
    "    (scores.get(r.id) ?? 0) + (1-bm25Weight)/(k+i+1)));\n"
    "  bm25R.forEach((r, i) => scores.set(r.id,\n"
    "    (scores.get(r.id) ?? 0) + bm25Weight/(k+i+1)));\n"
    "  return [...scores.entries()]\n"
    "    .sort((a,b) => b[1]-a[1]).slice(0,topK)\n"
    "    .map(([id]) => vecR.find(r=>r.id===id) ?? bm25R.find(r=>r.id===id)!);\n"
    "}"
), CX, CY, CW, Inches(4.2), size=Pt(8.5))

# ── Slide 50 ──────────────────────────────────────────────────────────────────
s = content_slide("Prompt Assembly", "Part 6 \u00b7 RAG Pipeline")
code_block(s, (
    "function buildRAGPrompt(query: string, chunks: Chunk[]) {\n"
    "  const context = chunks\n"
    "    .map((c, i) => `[${i+1}] ${c.text}`)\n"
    "    .join(\"\\n\\n\");\n"
    "  return [\n"
    "    { role: \"system\", content:\n"
    "      \"Answer using ONLY the context below.\\n\"\n"
    "      + \"If the answer is not in the context, say so.\\n\\n\"\n"
    "      + context\n"
    "    },\n"
    "    { role: \"user\", content: query }\n"
    "  ];\n"
    "}"
), CX, CY, Inches(5.5), Inches(3.0), size=Pt(9))
buls(s, ["\"ONLY the context\" reduces hallucination",
         "Number chunks so LLM can cite sources",
         "Keep context under 2048 tokens for Gemma 3 1B",
         "Prepend agent system prompt before RAG context"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# =============================================================================
# PART 7 — Tools
# =============================================================================
section_slide("Part 7\nTools", "Slides 51\u201354")

# ── Slide 51 ──────────────────────────────────────────────────────────────────
s = content_slide("Tool Definition Schema", "Part 7 \u00b7 Tools")
code_block(s, (
    "interface Tool {\n"
    "  name: string;\n"
    "  description: string;   // LLM reads this to decide when to call\n"
    "  parameters: {\n"
    "    type: \"object\";\n"
    "    properties: Record<string, {\n"
    "      type: string;\n"
    "      description: string;\n"
    "      enum?: string[];\n"
    "    }>;\n"
    "    required: string[];\n"
    "  };\n"
    "  execute(args: Record<string, unknown>): Promise<string>;\n"
    "}"
), CX, CY, Inches(5.5), Inches(3.2), size=Pt(9))
buls(s, ["description is the most important field \u2014 be specific",
         "execute() returns a string the LLM reads as tool output",
         "Tools are injected into the system prompt as JSON schema",
         "LLM emits <tool_call>{...}</tool_call> to invoke"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 52 ──────────────────────────────────────────────────────────────────
s = content_slide("Built-in Tools", "Part 7 \u00b7 Tools")
table(s, ["Tool", "Description", "Key parameters"],
      [["search_documents", "Hybrid RAG search over ingested docs", "query: string"],
       ["get_current_time", "Returns ISO timestamp", "(none)"],
       ["calculate",        "Evaluate a math expression",           "expression: string"],
       ["get_weather",      "Fetch weather for a city (demo)",      "city: string"],
       ["list_documents",   "List all ingested document titles",    "(none)"],
       ["web_search",       "DuckDuckGo instant answer (demo)",     "query: string"]],
      CX, CY, CW, Inches(3.2), size=Pt(11))
txb(s, "Tools are assigned per-agent in the agent config \u2014 not globally available",
    CX, Inches(4.4), CW, Inches(0.4), size=Pt(11), color=CB_GREY, italic=True)

# ── Slide 53 ──────────────────────────────────────────────────────────────────
s = content_slide("ReAct Loop Implementation", "Part 7 \u00b7 Tools")
code_block(s, (
    "async function* runReAct(messages, tools, maxIter = 5) {\n"
    "  const history = [...messages];\n"
    "  for (let iter = 0; iter < maxIter; iter++) {\n"
    "    const response = yield* streamLLM(history);\n"
    "    const calls = parseToolCalls(response);\n"
    "    if (calls.length === 0) return;          // done\n\n"
    "    history.push({ role: \"assistant\", content: response });\n"
    "    for (const call of calls) {\n"
    "      const tool = tools.find(t => t.name === call.name);\n"
    "      if (!tool) { history.push({ role: \"tool\",\n"
    "        content: `Unknown tool: ${call.name}` }); continue; }\n"
    "      const result = await tool.execute(call.args);\n"
    "      history.push({ role: \"tool\", content: result });\n"
    "    }\n"
    "  }\n"
    "  yield* streamLLM(history);  // final answer after max iterations\n"
    "}"
), CX, CY, CW, Inches(4.0), size=Pt(8.5))
txb(s, "yield* streams tokens to the UI while the loop runs",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_GREEN, bold=True)

# ── Slide 54 ──────────────────────────────────────────────────────────────────
s = content_slide("Adding a Custom Tool", "Part 7 \u00b7 Tools")
code_block(s, (
    "// src/lib/tools/myTool.ts\nexport const myTool: Tool = {\n"
    "  name: \"lookup_product\",\n"
    "  description: \"Look up a product by SKU and return its details.\",\n"
    "  parameters: {\n"
    "    type: \"object\",\n"
    "    properties: {\n"
    "      sku: { type: \"string\", description: \"Product SKU code\" }\n"
    "    },\n"
    "    required: [\"sku\"]\n"
    "  },\n"
    "  async execute({ sku }) {\n"
    "    const product = await fetchProduct(String(sku));\n"
    "    return JSON.stringify(product);\n"
    "  }\n"
    "};\n\n"
    "// Register in src/lib/tools/index.ts\nexport const ALL_TOOLS = [searchDocuments, getCurrentTime, myTool];"
), CX, CY, CW, Inches(4.0), size=Pt(8.5))
txb(s, "Assign to an agent: agent.tools = [\"lookup_product\"]",
    CX, Inches(4.75), CW, Inches(0.4), size=Pt(11), color=CB_GREEN, bold=True)

# =============================================================================
# PART 8 — Agents
# =============================================================================
section_slide("Part 8\nAgents", "Slides 55\u201358")

# ── Slide 55 ──────────────────────────────────────────────────────────────────
s = content_slide("Agent Data Model", "Part 8 \u00b7 Agents")
code_block(s, (
    "interface Agent {\n"
    "  id: string;\n"
    "  name: string;\n"
    "  description: string;    // shown to router LLM\n"
    "  systemPrompt: string;   // injected before every conversation\n"
    "  tools: string[];        // tool names from ALL_TOOLS\n"
    "  ragEnabled: boolean;    // inject retrieved chunks?\n"
    "  modelPath?: string;     // override default model\n"
    "  temperature?: number;   // 0\u20131, default 0.7\n"
    "  maxTokens?: number;     // default 512\n"
    "}"
), CX, CY, Inches(5.5), Inches(2.8), size=Pt(9))
buls(s, ["Agents stored in _default.agents collection",
         "Default agent created by migration v1",
         "UI: Settings \u2192 Agents \u2192 New Agent",
         "Router sees all agent names + descriptions"],
     Inches(5.8), CY + Inches(0.3), Inches(3.86), Inches(2.5), size=Pt(12))

# ── Slide 56 ──────────────────────────────────────────────────────────────────
s = content_slide("Agent Router", "Part 8 \u00b7 Agents")
code_block(s, (
    "async function routeMessage(message: string, agents: Agent[]) {\n"
    "  const prompt = [\n"
    "    { role: \"system\", content:\n"
    "      \"You are a router. Given a user message, return JSON:\\n\"\n"
    "      + '{\"agent\": \"<agent_name>\"}\\n'\n"
    "      + \"Choose from: \" + agents.map(a =>\n"
    "          `${a.name}: ${a.description}`).join(\"; \")\n"
    "    },\n"
    "    { role: \"user\", content: message }\n"
    "  ];\n"
    "  const response = await runLLM(prompt, { maxTokens: 20 });\n"
    "  const { agent } = JSON.parse(extractJSON(response));\n"
    "  return agents.find(a => a.name === agent) ?? agents[0];\n"
    "}"
), CX, CY, CW, Inches(3.5), size=Pt(9))
buls(s, ["maxTokens: 20 \u2014 router only needs a short JSON response",
         "extractJSON() strips markdown code fences if present",
         "Falls back to agents[0] (default) on parse error",
         "Skip routing if only one agent configured"],
     CX, Inches(4.3), CW, Inches(1.3), size=Pt(13))

# ── Slide 57 ──────────────────────────────────────────────────────────────────
s = content_slide("Conversation State Machine")
_bw = Inches(2.0); _bh = Inches(0.5); _gap = Inches(0.12); _arr = Inches(0.22)
_bx = CX + Inches(0.2); _ty = CY
_states = [
    ("idle",            RGBColor(0xE5,0xE5,0xE5), ""),
    ("routing",         RGBColor(0xFF,0xF0,0xDB), "router LLM call"),
    ("generating",      RGBColor(0xEC,0x12,0x18), "streaming tokens"),
    ("executing_tool",  RGBColor(0xF2,0xF2,0xF2), "tool result appended"),
    ("idle",            RGBColor(0xE5,0xE5,0xE5), "stream ends"),
]
for i,(lbl,fc,note) in enumerate(_states):
    tc = CB_WHITE if fc == RGBColor(0xEC,0x12,0x18) else CB_DARK
    diag_box(s, lbl, "", _bx, _ty, _bw, _bh, fc=fc, tc=tc, size=Pt(11))
    if note:
        txb(s, note, _bx+_bw+Inches(0.12), _ty, Inches(2.4), _bh,
            size=Pt(9), color=CB_GREY, italic=True)
    _ty += _bh
    if i < len(_states)-1:
        diag_label(s, "\u2193", _bx, _ty, _bw, _arr, color=CB_RED, size=Pt(13), bold=True)
        _ty += _arr
# Side branch: generating -> executing_tool
txb(s, "\u2192 tool_call detected", _bx+_bw+Inches(0.12),
    CY + 2*(_bh+_arr) + Inches(0.05), Inches(2.4), Inches(0.35),
    size=Pt(9), color=CB_RED, italic=True)
buls(s, ["State held in useChat hook (React)",
         "Abort button transitions generating \u2192 idle",
         "Error transitions any state \u2192 idle + toast",
         "Messages persisted to DB at each state transition"],
     Inches(5.0), CY + Inches(0.3), Inches(4.66), Inches(2.5), size=Pt(12))

# ── Slide 58 ──────────────────────────────────────────────────────────────────
s = content_slide("useChat Hook Architecture", "Part 8 \u00b7 Agents")
code_block(s, (
    "// src/hooks/useChat.ts  (simplified)\nexport function useChat() {\n"
    "  const [messages, setMessages] = useState<Message[]>([]);\n"
    "  const [status, setStatus] = useState<ChatStatus>(\"idle\");\n\n"
    "  async function sendMessage(text: string) {\n"
    "    setStatus(\"routing\");\n"
    "    const agent = await routeMessage(text, agents);\n"
    "    setStatus(\"generating\");\n"
    "    const userMsg = await saveMessage({ role:\"user\", content:text });\n"
    "    setMessages(m => [...m, userMsg]);\n"
    "    let assistantContent = \"\";\n"
    "    for await (const token of runReAct([...messages, userMsg],\n"
    "                                       agent.tools)) {\n"
    "      assistantContent += token;\n"
    "      setMessages(m => updateLast(m, assistantContent));\n"
    "    }\n"
    "    await saveMessage({ role:\"assistant\", content:assistantContent });\n"
    "    setStatus(\"idle\");\n"
    "  }\n"
    "  return { messages, status, sendMessage };\n"
    "}"
), CX, CY, CW, Inches(4.2), size=Pt(8.5))


# =============================================================================
# PART 9 — Putting It All Together
# =============================================================================
section_slide("Part 9\nPutting It All Together", "Slides 59\u201361")

# ── Slide 59 ──────────────────────────────────────────────────────────────────
s = content_slide("End-to-End Request Flow", "Part 9 \u00b7 Putting It All Together")
# Vertical pipeline: 6 steps, left column; annotation notes on right
_fw = Inches(2.6); _fh = Inches(0.48); _farr = Inches(0.18)
_fx = CX + Inches(0.1); _fy = CY
_steps = [
    ("User types message",      CB_LGREY,  CB_DARK,   ""),
    ("routeMessage()",          CB_CREAM,  CB_DARK,   "LLM call  ~200 ms"),
    ("hybridSearch()",          CB_CREAM,  CB_DARK,   "embed + SQL++ + BM25 + RRF"),
    ("buildRAGPrompt()",        CB_CREAM,  CB_DARK,   "system + context + history"),
    ("runReAct()",              CB_RED,    CB_WHITE,  "stream tokens, execute tools"),
    ("saveMessage() + render",  CB_LGREY,  CB_DARK,   "persist to CouchbaseLite"),
]
for i, (lbl, fc, tc, note) in enumerate(_steps):
    diag_box(s, lbl, "", _fx, _fy, _fw, _fh, fc=fc, tc=tc, size=Pt(10))
    if note:
        txb(s, note, _fx + _fw + Inches(0.14), _fy, Inches(3.0), _fh,
            size=Pt(9), color=CB_GREY, italic=True)
    _fy += _fh
    if i < len(_steps) - 1:
        diag_label(s, "\u2193", _fx, _fy, _fw, _farr, color=CB_RED, size=Pt(12), bold=True)
        _fy += _farr
buls(s, ["Total latency: ~200 ms routing + ~500 ms first token",
         "Streaming hides generation latency",
         "All data stays on-device \u2014 no cloud calls"],
     Inches(5.9), CY + Inches(0.4), Inches(3.76), Inches(2.0), size=Pt(12))

# ── Slide 60 ──────────────────────────────────────────────────────────────────
s = content_slide("Exercises", "Part 9 \u00b7 Putting It All Together")
table(s, ["#", "Exercise", "Difficulty"],
      [["1", "Add a new tool: fetch a URL and return its text",          "\u2605\u2605\u2606\u2606\u2606"],
       ["2", "Create a second agent with a different system prompt",      "\u2605\u2606\u2606\u2606\u2606"],
       ["3", "Ingest a PDF and ask questions about it",                   "\u2605\u2606\u2606\u2606\u2606"],
       ["4", "Tune bm25Weight and compare retrieval quality",             "\u2605\u2605\u2606\u2606\u2606"],
       ["5", "Add a Couchbase Sync Gateway and sync to a second device",  "\u2605\u2605\u2605\u2606\u2606"],
       ["6", "Replace Gemma 3 1B with a different .litertlm model",       "\u2605\u2605\u2606\u2606\u2606"],
       ["7", "Add a conversation export tool (Markdown / JSON)",          "\u2605\u2605\u2606\u2606\u2606"],
       ["8", "Implement semantic chunking (split on paragraphs)",         "\u2605\u2605\u2605\u2606\u2606"]],
      CX, CY, CW, Inches(4.0), size=Pt(11))

# ── Slide 61 ──────────────────────────────────────────────────────────────────
s = content_slide("Resources & Next Steps", "Part 9 \u00b7 Putting It All Together")
table(s, ["Resource", "URL"],
      [["This repo",          "github.com/ldoguin/tauri-cblite-litert"],
       ["Tauri docs",         "tauri.app/v2"],
       ["Couchbase Lite",     "docs.couchbase.com/couchbase-lite"],
       ["LiteRT / GenAI API", "ai.google.dev/edge/litert"],
       ["Gemma models",       "ai.google.dev/gemma"],
       ["BERT embeddings",    "huggingface.co/bert-base-uncased"],
       ["python-pptx",        "python-pptx.readthedocs.io"]],
      CX, CY, Inches(9.32), Inches(3.5), size=Pt(12))
buls(s, ["Star the repo \u2014 PRs welcome!",
         "Questions? Open a GitHub Discussion",
         "Workshop slides generated by make_pptx.py \u2014 edit and re-run"],
     CX, Inches(4.6), CW, Inches(1.0), size=Pt(14))

# =============================================================================
# APPENDIX — Glossary
# =============================================================================
section_slide("Appendix\nGlossary", "Slides 62\u201363")

# ── Slide 62 ── Glossary: AI & ML Terms ──────────────────────────────────────
s = content_slide("Glossary: AI & ML Terms")
table(s, ["Term", "Stands for / Meaning"],
      [
       ["LLM",      "Large Language Model \u2014 a transformer trained to predict the next token"],
       ["BERT",     "Bidirectional Encoder Representations from Transformers \u2014 encoder-only model for embeddings"],
       ["RAG",      "Retrieval-Augmented Generation \u2014 inject retrieved context into an LLM prompt"],
       ["KV-cache", "Key-Value cache \u2014 stores attention K/V tensors to avoid recomputing previous tokens"],
       ["INT4/INT8","4-bit / 8-bit integer quantisation \u2014 reduces model size and speeds up inference"],
       ["FP16/FP32","16-bit / 32-bit floating point \u2014 FP32 is full precision, FP16 halves memory"],
       ["PTQ",      "Post-Training Quantisation \u2014 quantise a trained model without retraining"],
       ["QAT",      "Quantisation-Aware Training \u2014 simulate quantisation during fine-tuning for better accuracy"],
       ["FFN",      "Feed-Forward Network \u2014 the two-layer MLP inside each transformer block"],
       ["MLM",      "Masked Language Modelling \u2014 BERT pre-training task: predict masked tokens"],
       ["NSP",      "Next Sentence Prediction \u2014 BERT pre-training task: predict if two sentences are adjacent"],
       ["SBERT",    "Sentence-BERT \u2014 BERT fine-tuned with a pooling head for sentence-level similarity"],
       ["BM25",     "Best Match 25 \u2014 probabilistic keyword ranking function (Okapi BM25)"],
       ["RRF",      "Reciprocal Rank Fusion \u2014 merges ranked lists from multiple retrieval methods"],
       ["IDF",      "Inverse Document Frequency \u2014 down-weights terms that appear in many documents"],
       ["GELU",     "Gaussian Error Linear Unit \u2014 smooth activation function used in transformers"],
       ["EOS",      "End-of-Sequence token \u2014 special token that signals the model to stop generating"],
       ["CLS",      "[CLS] token \u2014 BERT\u2019s first input token; its output vector represents the whole sequence"],
       ["SEP",      "[SEP] token \u2014 BERT separator between sentence pairs; also marks end of input"],
       ["PAD",      "[PAD] token \u2014 pads sequences to a fixed length; ignored by the attention mask"],
      ],
      CX, CY, CW, Inches(5.1), size=Pt(9))

# ── Slide 63 ── Glossary: Infrastructure & Format Terms ──────────────────────
s = content_slide("Glossary: Infrastructure & Format Terms")
table(s, ["Term", "Stands for / Meaning"],
      [
       ["IPC",      "Inter-Process Communication \u2014 how the Tauri frontend (JS) calls the backend (Rust)"],
       ["ACL",      "Access Control List \u2014 Tauri v2 permission system that gates which commands the frontend can call"],
       ["GGUF",     "GPT-Generated Unified Format \u2014 llama.cpp\u2019s binary model format (successor to GGML)"],
       ["ONNX",     "Open Neural Network Exchange \u2014 open format for representing ML models across frameworks"],
       ["GPU",      "Graphics Processing Unit \u2014 massively parallel processor; accelerates matrix operations in ML"],
       ["CPU",      "Central Processing Unit \u2014 general-purpose processor; fallback when no GPU delegate is available"],
       ["NPU",      "Neural Processing Unit \u2014 dedicated silicon for ML inference (e.g. Qualcomm Hexagon, Apple ANE)"],
       ["ANE",      "Apple Neural Engine \u2014 Apple\u2019s on-chip NPU in A-series and M-series chips"],
       ["DSP",      "Digital Signal Processor \u2014 low-power fixed-function chip; used by Qualcomm NNAPI delegate"],
       ["NNAPI",    "Neural Networks API \u2014 Android abstraction layer that routes inference to GPU/DSP/NPU"],
       ["SIMD",     "Single Instruction Multiple Data \u2014 CPU vector instructions (NEON on ARM, AVX on x86)"],
       ["ACID",     "Atomicity, Consistency, Isolation, Durability \u2014 database transaction guarantees"],
       ["MVCC",     "Multi-Version Concurrency Control \u2014 allows concurrent reads without blocking writes"],
       ["SQL++",    "Superset of SQL for JSON documents \u2014 Couchbase\u2019s query language (also called N1QL)"],
       ["CDN",      "Content Delivery Network \u2014 geographically distributed servers for fast file downloads"],
       ["OTA",      "Over-The-Air \u2014 delivering updates (models, config) without an app store release"],
       ["SHA",      "Secure Hash Algorithm \u2014 SHA-256 used here to verify model file integrity after download"],
       ["FFI",      "Foreign Function Interface \u2014 mechanism for calling native C/Rust code from another language"],
       ["SDK",      "Software Development Kit \u2014 bundled libraries, tools, and docs for a platform or API"],
       ["HNSW",     "Hierarchical Navigable Small World \u2014 graph-based approximate nearest-neighbour index"],
      ],
      CX, CY, CW, Inches(5.1), size=Pt(9))

# =============================================================================
# Save
# =============================================================================
prs.save("workshop.pptx")
print(f"Saved workshop.pptx  ({len(prs.slides)} slides)")
